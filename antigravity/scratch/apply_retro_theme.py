import pptx
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

def set_shape_text_style(shape, font_name='Courier New', font_color=RGBColor(0, 255, 0), is_title=False):
    if not shape.has_text_frame:
        return

    text_frame = shape.text_frame
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.color.rgb = font_color
            if is_title:
                run.font.bold = True
                # Increase size for title if needed?
                # run.font.size = Pt(40) 

def apply_theme(input_path, output_path):
    print(f"Loading: {input_path}")
    prs = pptx.Presentation(input_path)

    # Define Colors
    bg_color = RGBColor(0, 0, 0)       # Black
    text_color = RGBColor(0, 255, 0)   # Green (Terminal)
    title_font = 'Courier New'
    body_font = 'Courier New'

    # Apply to Masters (Layouts)
    # Be careful, modifying masters might break specific formatting, but needed for consistency.
    print("Applying theme to Masters...")
    for slide_master in prs.slide_masters:
        # Set background
        bg = slide_master.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
        
        for layout in slide_master.slide_layouts:
            bg = layout.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = bg_color
            
            # Change placeholders on layout
            for shape in layout.placeholders:
                is_title = (shape.name and 'Title' in shape.name)
                set_shape_text_style(shape, title_font, text_color, is_title)

    # Apply to Slides
    print("Applying theme to Slides...")
    for slide in prs.slides:
        # Background again (some slides override master)
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
        
        for shape in slide.shapes:
            # Check type
            if shape.has_text_frame:
                is_title = (shape == slide.shapes.title)
                set_shape_text_style(shape, body_font, text_color, is_title)
            
            # Also handle groups?
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for subshape in shape.shapes:
                    if subshape.has_text_frame:
                         set_shape_text_style(subshape, body_font, text_color)
            
            # Handle Tables?
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text_frame:
                            set_shape_text_style(cell, body_font, text_color)

    prs.save(output_path)
    print(f"Saved themed presentation to: {output_path}")

input_path = r"C:\Users\Administrator\Downloads\Operating_System_Design_Engineering_Final_Organized.pptx"
output_path = r"C:\Users\Administrator\Downloads\Operating_System_Design_Engineering_Final_Organized_Themed.pptx"

if __name__ == "__main__":
    apply_theme(input_path, output_path)
