from pptx import Presentation
from pptx.util import Inches, Pt
import os

path_theme = r"c:\Users\Administrator\Downloads\Lecture_Grp_D.pptx"
path_content = r"c:\Users\Administrator\Downloads\Operating_System_Design_Merged.pptx"
output_path = r"c:\Users\Administrator\Downloads\Operating_System_Design_User_Structure_Themed.pptx"

print("Loading presentations...")
# We use the theme file as the base to get the masters/layouts
prs_template = Presentation(path_theme) 
prs_content = Presentation(path_content)

# Clear all slides in template
xml_slides = prs_template.slides._sldIdLst
slides_to_remove = list(xml_slides)
for s in slides_to_remove:
    xml_slides.remove(s)

print("Applying theme to user slides...")

def copy_content(source, target):
    # Title
    if source.shapes.title and target.shapes.title:
        target.shapes.title.text = source.shapes.title.text
        
    # Body
    src_body = None
    for s in source.placeholders:
        if s.placeholder_format.idx == 1:
            src_body = s
            break
            
    tgt_body = None
    for s in target.placeholders:
        if s.placeholder_format.idx == 1:
            tgt_body = s
            break
            
    if src_body and src_body.has_text_frame and tgt_body:
        tgt_body.text_frame.clear()
        for p in src_body.text_frame.paragraphs:
            new_p = tgt_body.text_frame.add_paragraph()
            new_p.text = p.text
            new_p.level = p.level
            # Preserve bold/italic but let theme dictate font/color
            if p.font.bold: new_p.font.bold = True
            if p.font.italic: new_p.font.italic = True

    # Images
    for shape in source.shapes:
        if shape.shape_type == 13: # PICTURE
             try:
                image = shape.image
                filename = f"temp_strict_{shape.shape_id}.{image.ext}"
                with open(filename, 'wb') as f:
                    f.write(image.blob)
                target.shapes.add_picture(filename, shape.left, shape.top, shape.width, shape.height)
                os.remove(filename)
             except: pass

# Iterate User Slides and append to Template
for i, slide in enumerate(prs_content.slides):
    # Choose Layout
    # Mapping:
    # Title Slide -> 0
    # Section Header -> 1
    # Title Only -> 4 (usually)
    # Content -> 2
    
    name = slide.slide_layout.name
    layout_idx = 2
    if "Title Slide" in name: layout_idx = 0
    elif "Section Header" in name: layout_idx = 1
    elif "Title Only" in name: layout_idx = 4
    
    if layout_idx >= len(prs_template.slide_layouts): layout_idx = 2
    
    layout = prs_template.slide_layouts[layout_idx]
    new_slide = prs_template.slides.add_slide(layout)
    
    copy_content(slide, new_slide)

print(f"Saving to {output_path}...")
prs_template.save(output_path)
print("Done.")
