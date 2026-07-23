import pptx
import os

def move_slide(prs, old_index, new_index):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[old_index])
    xml_slides.insert(new_index, slides[old_index])

def reorder_presentation(filepath, output_path):
    print(f"Loading: {filepath}")
    prs = pptx.Presentation(filepath)
    
    # 1. Identify the 'Quiz' or 'Summary' start index
    split_index = -1
    for i, slide in enumerate(prs.slides):
        if slide.shapes.title:
            title = slide.shapes.title.text.lower()
            if "quiz" in title or "summary" in title:
                split_index = i
                print(f"Found split point at index {i}: {title}")
                break
    
    if split_index == -1:
        print("Could not find 'Quiz' or 'Summary' section. No re-ordering needed.")
        return

    # 2. Identify the new appended slides
    # We assume everything AFTER the original 'Q & A' (original last slide) is new.
    # But wait, we appended them to the end.
    # Original length was 34? Let's check the logs.
    # From previous 'analyze_pptx' on target: Slide 34: Q & A.  So original length was 34 slides (indices 0-33).
    # New slides are at indices 34 to End.
    
    # We want to move everything from 34 to End -> into split_index.
    
    # Get total count
    total_slides = len(prs.slides)
    original_end_index = 34 # 0-indexed, this is the index of the first NEW slide
    
    # Verify assumption
    if total_slides <= original_end_index:
        print("No new slides found to reorder.")
        return

    print(f"Moving slides {original_end_index} to {total_slides-1} to position {split_index}")

    # We move them one by one.
    # When we move a slide from 'original_end_index' to 'split_index', 
    # the 'original_end_index' effectively increments for the NEXT slide we want to move 
    # because we shifted things down? NO.
    # Actually, it's safer to move them one by one.
    
    # Example: [A, B, C, D, E] -> Move D, E before C.
    # Move D to 2: [A, B, D, C, E]
    # Move E to 3: [A, B, D, E, C]
    
    # So we iterate through the range of new slides.
    # Since we are modifying the list, let's be careful.
    # The 'new' slides are currently at [34, 35, ... 52]
    # We want to move slide 34 to 'split_index' (29)
    # Then slide 35 (which is now at 35) to 'split_index + 1' (30)
    
    num_new_slides = total_slides - original_end_index
    current_new_slide_index = original_end_index
    target_insert_index = split_index
    
    for _ in range(num_new_slides):
        move_slide(prs, current_new_slide_index, target_insert_index)
        # After moving 34 to 29. The old 29 becomes 30. The old 34 is now at 29.
        # The slide that WAS at 35 is now at 35 still? 
        # Wait. 
        # [0...28] [29: Quiz] ... [33: Q&A] [34: New1] [35: New2]
        # Move 34 to 29:
        # [0...28] [29: New1] [30: Quiz] ... [34: Q&A] [35: New2]
        # Now we want to move New2 (at 35) to 30.
        # So 'current_new_slide_index' stays the same?
        # Yes, because the 'hole' left by moving 34 is filled by shifting 35 down? 
        # No, 'slides' list shifts right.
        # Let's trace list operations.
        # remove(34). insert(29).
        # List becomes: 0..28, 34(New1), 29(Quiz)...33, 35(New2)... 
        # The slide at index 35 is now New2.
        # So yes, we can keep grabbing from 35?
        # Wait, if we grab from 35 and insert at 30.
        # remove(35). insert(30).
        # List: 0..28, New1, New2, Quiz... 
        # So we always grab from 'original_end_index + i' ?
        # Actually, let's just observe that after moving 34 to 29, 
        # the slide that WAS at 35 moves to 35? No, indices shift.
        # Initial: 0..28, 29..33, 34..52
        # Move 34 to 29.
        # New Indices: 0..28, 29(New1), 30(Old29)..34(Old33), 35(New2)..53
        # The slide New2 is at 35. 
        # So yes, we just keep incrementing both.
        
        current_new_slide_index += 1
        target_insert_index += 1

    prs.save(output_path)
    print(f"Saved reordered presentation to: {output_path}")

merged_path = r"C:\Users\Administrator\Downloads\Operating_System_Design_Engineering_Final_Merged.pptx"
final_path = r"C:\Users\Administrator\Downloads\Operating_System_Design_Engineering_Final_Merged.pptx" # Overwrite? Better to verify first. 
# Let's overwrite for simplicity as per user flow, but maybe keep a backup?
# The user script logic above assumes 'merged' is the input.
# Let's output to the same file.

if __name__ == "__main__":
    reorder_presentation(merged_path, final_path)
