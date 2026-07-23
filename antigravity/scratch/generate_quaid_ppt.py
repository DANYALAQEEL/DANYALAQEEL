from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import json
import os

# --- Configuration ---
CONTENT_FILE = r"C:\Users\Administrator\.gemini\antigravity\scratch\quaid_content.json"
ASSETS_DIR = r"C:\Users\Administrator\.gemini\antigravity\scratch\assets"
OUTPUT_FILE = r"C:\Users\Administrator\Downloads\Speeches_of_Quaid_Redesigned.pptx"

# Colors
EMERALD_GREEN = RGBColor(0, 66, 37)
ANTIQUE_GOLD = RGBColor(184, 134, 11)
CREAM_WHITE = RGBColor(245, 245, 220) # Lighter cream for text
CHARCOAL = RGBColor(51, 51, 51)
WHITE = RGBColor(255, 255, 255)

# --- Helper Functions ---

def apply_text_style(shape, font_family="Garamond", font_size=Pt(18), color=CHARCOAL, bold=False, italic=False):
    if not shape.has_text_frame:
        return
    text_frame = shape.text_frame
    for paragraph in text_frame.paragraphs:
        paragraph.font.name = font_family
        paragraph.font.size = font_size
        paragraph.font.color.rgb = color
        paragraph.font.bold = bold
        paragraph.font.italic = italic

def add_background_image(slide, image_filename):
    img_path = os.path.join(ASSETS_DIR, image_filename)
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, 0, 0, width=Inches(13.333), height=Inches(7.5))

def create_title_slide(prs, slide_data):
    slide_layout = prs.slide_layouts[6] # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background: Emerald Green
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = EMERALD_GREEN
    bg.line.fill.background() # No line

    # Add Gold Frame
    add_background_image(slide, "gold_frame.png")
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = slide_data['title'] if slide_data['title'] else slide_data['body_text'][0]
    p.font.name = "Garamond"
    p.font.size = Pt(54)
    p.font.color.rgb = WHITE # Antique Gold
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if len(slide_data['body_text']) > 1:
        sub_box = slide.shapes.add_textbox(Inches(2), Inches(3.8), Inches(9.333), Inches(1))
        tf = sub_box.text_frame
        p = tf.add_paragraph()
        p.text = slide_data['body_text'][1]
        p.font.name = "Calibri Light"
        p.font.size = Pt(28)
        p.font.color.rgb = RGBColor(200, 200, 200)
        p.alignment = PP_ALIGN.CENTER

    # Footer
    if len(slide_data['body_text']) > 2:
        foot_box = slide.shapes.add_textbox(Inches(2), Inches(6), Inches(9.333), Inches(1))
        tf = foot_box.text_frame
        p = tf.add_paragraph()
        p.text = slide_data['body_text'][2]
        p.font.name = "Calibri"
        p.font.size = Pt(14)
        p.font.color.rgb = ANTIQUE_GOLD
        p.alignment = PP_ALIGN.CENTER

def create_standard_slide(prs, slide_data, layout_type="standard"):
    slide_layout = prs.slide_layouts[6] # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background: Texture
    add_background_image(slide, "parchment_bg.png")
    
    # Header Bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), Inches(12.333), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = EMERALD_GREEN
    header.line.color.rgb = ANTIQUE_GOLD
    header.line.width = Pt(2)
    
    # Title Text
    tf = header.text_frame
    p = tf.paragraphs[0]
    # Check if title is empty, use first body line if so (common in extraction)
    title_text = slide_data['title']
    body_start_idx = 0
    if not title_text and slide_data['body_text']:
         title_text = slide_data['body_text'][0]
         body_start_idx = 1
         
    p.text = title_text
    p.font.name = "Garamond"
    p.font.size = Pt(36)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    header.text_frame.margin_left = Inches(0.2)
    
    # Content Area
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    body_content = slide_data['body_text'][body_start_idx:]
    
    for line in body_content:
        # Detect sub-headings or emphasis based on length or capitalization
        p = tf.add_paragraph()
        p.text = line
        p.font.name = "Calibri"
        p.space_after = Pt(10)
        
        if len(line) < 50 and line[0].isupper() and ":" not in line:
            # Likely a sub-header
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = EMERALD_GREEN
        else:
            # Normal text
            p.font.size = Pt(20)
            p.font.color.rgb = CHARCOAL

def create_quote_slide(prs, slide_data):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Dark Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(20, 30, 20) # Very dark green
    
    # Decorative Pattern
    # Just place the pattern in corners
    pattern_path = os.path.join(ASSETS_DIR, "pattern.png")
    if os.path.exists(pattern_path):
        slide.shapes.add_picture(pattern_path, Inches(0.5), Inches(0.5), width=Inches(2), height=Inches(2))
        slide.shapes.add_picture(pattern_path, Inches(10.8), Inches(5), width=Inches(2), height=Inches(2))

    # Quote Text
    quote_text = "\n".join(slide_data['body_text'][1:]) # Assuming index 0 is title
    
    textbox = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(9.333), Inches(4))
    tf = textbox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = f'"{quote_text}"'
    p.font.name = "Garamond"
    p.font.size = Pt(32)
    p.font.italic = True
    p.font.color.rgb = RGBColor(240, 230, 140) # Khaki/Gold text
    p.alignment = PP_ALIGN.CENTER

    # Source styling
    source_p = tf.add_paragraph()
    source_p.text = "- Quaid-e-Azam Muhammad Ali Jinnah"
    source_p.font.name = "Calibri"
    source_p.font.size = Pt(20)
    source_p.font.color.rgb = WHITE
    source_p.alignment = PP_ALIGN.RIGHT
    source_p.space_before = Pt(20)

def main():
    if not os.path.exists(CONTENT_FILE):
        print("Content file not found.")
        return

    with open(CONTENT_FILE, "r", encoding="utf-8") as f:
        data = json.load(f)

    prs = Presentation()
    # 16:9 aspect ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides = data['slides']
    
    # --- Generation Loop ---
    for i, slide_data in enumerate(slides):
        print(f"Generating slide {i}: {slide_data.get('title', 'Untitled')}")
        
        if i == 0:
            create_title_slide(prs, slide_data)
        elif "11 August 1947" in str(slide_data.get('body_text')) or "Lahore Resolution" in str(slide_data.get('body_text')):
            # treat key speeches as quote/feature slides if text is short enough, otherwise standard
            # Check length - if body text is huge, standard is safer
            text_len = sum(len(s) for s in slide_data['body_text'])
            if text_len < 400 and ("11 August" in str(slide_data.get('body_text'))):
                 create_quote_slide(prs, slide_data)
            else:
                 create_standard_slide(prs, slide_data)
        else:
            create_standard_slide(prs, slide_data)

    prs.save(OUTPUT_FILE)
    print(f"Presentation saved to {OUTPUT_FILE}")

if __name__ == "__main__":
    main()
