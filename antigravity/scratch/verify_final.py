from pptx import Presentation
import os

path = r"c:\Users\Administrator\Downloads\Operating_System_Design_Final_Dark_Fixed.pptx"

if not os.path.exists(path):
    print("File not found.")
    exit()

try:
    prs = Presentation(path)
    print(f"Presentation loaded successfully. Total slides: {len(prs.slides)}")
    
    print("\n--- Slide Titles ---")
    for i, slide in enumerate(prs.slides):
        title = slide.shapes.title.text if slide.shapes.title else "(No Title)"
        print(f"{i+1}. {title}")
        
except Exception as e:
    print(f"Error loading presentation: {e}")
