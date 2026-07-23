import pptx
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

def set_text_frame_style(text_frame, font_name='Courier New', font_color=RGBColor(0, 255, 0), is_title=False):
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.color.rgb = font_color
            if is_title:
                run.font.bold = True

def process_shape(shape, font_name, font_color, is_title=False):
    if shape.has_text_frame:
        set_text_frame_style(shape.text_frame, font_name, font_color, is_title)
    
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                if cell.text_frame:
                    set_text_frame_style(cell.text_frame, font_name, font_color)
    
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for subshape in shape.shapes:
            process_shape(subshape, font_name, font_color)

def apply_theme(input_path, output_path):
    print(f"Loading: {input_path}")
    prs = pptx.Presentation(input_path)

    bg_color = RGBColor(0, 0, 0)       # Black
    text_color = RGBColor(0, 255, 0)   # Green
    font_name = 'Courier New'

    # 1. Masters
    print("Applying to Masters...")
    for master in prs.slide_masters:
        master.background.fill.solid()
        master.background.fill.fore_color.rgb = bg_color
        # Check text styles in master?
        # python-pptx limitation: can't easily set default theme fonts.
        # But we can iterate layouts.
        for layout in master.slide_layouts:
            layout.background.fill.solid()
            layout.background.fill.fore_color.rgb = bg_color
            for shape in layout.placeholders:
                is_title = (shape.name and 'Title' in shape.name)
                process_shape(shape, font_name, text_color, is_title)
            for shape in layout.shapes:
                 process_shape(shape, font_name, text_color)

    # 2. Slides
    print("Applying to Slides...")
    for slide in prs.slides:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = bg_color
        
        is_title_slide = False
        if slide.shapes.title:
            is_title_slide = True
            
        for shape in slide.shapes:
            is_title = (shape == slide.shapes.title)
            process_shape(shape, font_name, text_color, is_title)

    prs.save(output_path)
    print(f"Saved fixed presentation to: {output_path}")

input_path = r"C:\Users\Administrator\Downloads\Operating_System_Design_Engineering_Final_Organized.pptx" # Start from organized
output_path = r"C:\Users\Administrator\Downloads\Operating_System_Design_Engineering_Final_Organized_Themed_Fixed.pptx"

if __name__ == "__main__":
    apply_theme(input_path, output_path)
