from pptx import Presentation
import os

path = r"c:\Users\Administrator\Downloads\Operating_System_Design_User_Structure_Themed.pptx"

if not os.path.exists(path):
    print("File not found.")
    exit()

try:
    prs = Presentation(path)
    print(f"Loaded: {len(prs.slides)} slides.")
    # Print first 5 titles to confirm structure
    for i, slide in enumerate(prs.slides[:5]):
       t = slide.shapes.title.text if slide.shapes.title else "(No Title)"
       print(f"{i+1}. {t}")
except Exception as e:
    print(f"Error: {e}")
