from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
import copy
import os

# Paths
path_counter = r"c:\Users\Administrator\Downloads\GroupD-SectionA-Topic-2.7.1-2.7.3 (1).pptx"
path_user = r"c:\Users\Administrator\Downloads\Operating_System_Design_Merged.pptx"
output_path = r"c:\Users\Administrator\Downloads\Operating_System_Design_Final_Dark.pptx"

# Theme Colors (Dark Mode "Old Computer")
BG_COLOR = RGBColor(10, 10, 10)       # Near Black
TEXT_TITLE = RGBColor(0, 255, 65)     # Matrix Green / Retro Terminal
TEXT_BODY = RGBColor(220, 220, 220)   # Off-White
ACCENT_COLOR = RGBColor(0, 128, 0)    # Dark Green for lines

# Load Presentations
print("Loading presentations...")
prs_counter = Presentation(path_counter)
prs_user = Presentation(path_user)

# Helper: Normalize text for matching
def normalize(text):
    return text.lower().strip().replace(":", "").replace("-", " ")

# Index User Slides by Title
user_map = {}
for i, slide in enumerate(prs_user.slides):
    title = slide.shapes.title.text if slide.shapes.title else ""
    norm_title = normalize(title)
    if norm_title:
        if norm_title not in user_map:
            user_map[norm_title] = []
        user_map[norm_title].append(slide)

# Create Output Presentation (Start Fresh or from Template?)
# We'll start from User's base to keep layouts, but clear slides.
prs_out = Presentation(path_user)
xml_slides = prs_out.slides._sldIdLst
# Clear all existing slides to rebuild order
for i in range(len(xml_slides)):
    if len(xml_slides) > 0:
        xml_slides.remove(xml_slides[0])

# Function to apply theme to a slide
def apply_theme(slide):
    # 1. Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR
    
    # 2. Text styling
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                # Title
                if shape == slide.shapes.title:
                    p.font.color.rgb = TEXT_TITLE
                    p.font.name = "Consolas" # Monospaced code-like
                    p.font.bold = True
                else:
                    # Body
                    font = p.font
                    font.color.rgb = TEXT_BODY
                    font.name = "Consolas" # Monospaced code-like
                    
        # Invert specific shapes logic if needed?
        # Only if shape has fill.
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            if shape.fill.type:
                 # Check if it's a solid fill
                 try:
                     # If it's very light, make it dark grey
                     # This is hard to detect without reading RGB, skipping for safety
                     pass
                 except: pass

# Function to copy slide content to target presentation
def copy_slide(source_slide, target_prs):
    # This is tricky because copying slides perfectly requires XML cloning
    # or complex shape iteration.
    # Python-pptx doesn't support full slide copy easily.
    # We will use the "add_slide" with similar layout and copy shapes.
    
    # Find matching layout
    layout_name = source_slide.slide_layout.name
    layout = target_prs.slide_layouts[0] # Default
    for l in target_prs.slide_layouts:
        if l.name == layout_name:
            layout = l
            break
            
    new_slide = target_prs.slides.add_slide(layout)
    
    # Copy Title
    if source_slide.shapes.title and new_slide.shapes.title:
        new_slide.shapes.title.text = source_slide.shapes.title.text
        
    # Copy Body Text
    # Heuristic: Find first body placeholder in source -> copy to first in target
    source_body = None
    for shape in source_slide.placeholders:
        if shape.placeholder_format.idx == 1:
            source_body = shape
            break
            
    if source_body and source_body.has_text_frame:
        target_body = None
        for shape in new_slide.placeholders:
            if shape.placeholder_format.idx == 1:
                target_body = shape
                break
        
        if target_body:
            target_body.text_frame.clear()
            for p in source_body.text_frame.paragraphs:
                new_p = target_body.text_frame.add_paragraph()
                new_p.text = p.text
                new_p.level = p.level
    
    # Copy Extra Text Boxes / Shapes (Simple approximation)
    # Copy Images
    for shape in source_slide.shapes:
        if shape.shape_type == 13: # PICTURE
             try:
                image = shape.image
                blob = image.blob
                filename = f"temp_{shape.shape_id}.{image.ext}"
                with open(filename, 'wb') as f:
                    f.write(blob)
                new_slide.shapes.add_picture(filename, shape.left, shape.top, shape.width, shape.height)
                os.remove(filename)
             except: pass
    
    apply_theme(new_slide)
    return new_slide

# --- EXECUTION: Reorder & Merge ---
print("Reordering based on Counter Group...")
used_user_titles = set()

for i, c_slide in enumerate(prs_counter.slides):
    c_title = c_slide.shapes.title.text if c_slide.shapes.title else ""
    norm_c_title = normalize(c_title)
    
    # Check if User has this topic (fuzzy match?)
    # or exact match
    
    print(f"Processing Topic: {c_title}")
    
    matched = False
    
    # Try exact/contains match in user map
    found_key = None
    for u_key in user_map:
        if u_key in norm_c_title or norm_c_title in u_key: # Bidirectional containment
            found_key = u_key
            break
            
    if found_key:
        print(f"  -> Found in User PPT: {found_key}")
        # Copy User Slide(s)
        for u_slide in user_map[found_key]:
            copy_slide(u_slide, prs_out)
            used_user_titles.add(found_key)
        matched = True
    else:
        # User misses this topic -> Copy Counter Group's slide
        print(f"  -> Missing in User PPT. Importing from Counter Group.")
        copy_slide(c_slide, prs_out) # Use Counter slide as source

# Append Remaining User Slides (orphans)
print("Appending remaining User slides...")
for u_key, u_slides in user_map.items():
    if u_key not in used_user_titles:
        print(f"  -> Appending extra: {u_key}")
        for u_slide in u_slides:
            copy_slide(u_slide, prs_out)

print(f"Saving to {output_path}...")
prs_out.save(output_path)
print("Done.")
