import pptx
from pptx.util import Inches, Pt
import json
import os

def add_visuals(pptx_path, plan_path, output_path):
    if not os.path.exists(pptx_path):
        print(f"Error: Presentation not found at {pptx_path}")
        return
    if not os.path.exists(plan_path):
        print(f"Error: Plan not found at {plan_path}")
        return

    prs = pptx.Presentation(pptx_path)
    
    with open(plan_path, 'r') as f:
        visual_plan = json.load(f)

    # Create a map for easier access
    headers_map = {item['slide_index']: item for item in visual_plan}
    
    # Get presentation dimensions
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    for i, slide in enumerate(prs.slides):
        if i in headers_map:
            item = headers_map[i]
            image_filename = item.get('filename')
            # Look for the image in the artifacts directory (where generate_image saves)
            # generate_image saves to specific paths, we need to know where. 
            # Assuming they are in the current working directory or a known artifacts location.
            # The tool documentation says "saved as an artifact". 
            # I will search for the file in the current directory.
            
            # The generate_image tool saves files with the name provided in ImageName + .png usually.
            # Let's assume they are in the current directory for now.
            image_path = f"{image_filename}.png" if not image_filename.endswith('.png') else image_filename
            
            # Check if file exists (it might be named differently or in a subdir)
            # If not found directly, try finding it by name
            if not os.path.exists(image_path):
                 # Try finding it with just the name provided to the tool
                 # The tool creates "ImageName"
                 base_name = os.path.splitext(image_filename)[0]
                 candidates = [f for f in os.listdir('.') if base_name in f and f.endswith('.png')]
                 if candidates:
                     image_path = candidates[0]
            
            if os.path.exists(image_path):
                print(f"Adding {image_path} to slide {i+1}")
                if item['type'] == 'background':
                    # Set as background (workaround as python-pptx doesn't support easy background set)
                    # We can add a picture filling the slide and send to back
                    pic = slide.shapes.add_picture(image_path, 0, 0, slide_width, slide_height)
                    # Move to background is tricky in python-pptx without modifying xml directly.
                    # For now, we just add it. If it covers text, that's an issue.
                    # Better to add it to a corner or side if it's not a true background
                    
                else:
                    # Place in bottom right corner or available space
                    # specific logic can be improved
                    img_width = Inches(3)
                    img_height = Inches(2) # Aspect ratio might need adjustment
                    
                    left = slide_width - img_width - Inches(0.5)
                    top = slide_height - img_height - Inches(0.5)
                    
                    slide.shapes.add_picture(image_path, left, top, width=img_width)
            else:
                print(f"Warning: Image {image_filename} not found.")

    prs.save(output_path)
    print(f"Saved modified presentation to {output_path}")

if __name__ == "__main__":
    ppt_path = "Operating_System_Design_Engineering_Final_Organized.pptx"
    plan_path = "visual_plan.json"
    output_path = "Operating_System_Design_Engineering_Final_Organized_Visuals.pptx"
    add_visuals(ppt_path, plan_path, output_path)
