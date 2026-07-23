import pptx
import os
import json

def extract_content(pptx_path, output_json_path):
    if not os.path.exists(pptx_path):
        print(f"Error: File not found at {pptx_path}")
        return

    try:
        prs = pptx.Presentation(pptx_path)
    except Exception as e:
        print(f"Error opening presentation: {e}")
        return

    content = {
        "presentation_title": os.path.basename(pptx_path),
        "slides": []
    }

    print(f"Analyzing {len(prs.slides)} slides...")

    for i, slide in enumerate(prs.slides):
        slide_data = {
            "index": i,
            "title": "",
            "body_text": [],
            "notes": ""
        }

        # Extract Title
        if slide.shapes.title:
            slide_data["title"] = slide.shapes.title.text.strip()

        # Extract Body Text
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text = shape.text.strip()
                # Avoid duplicating title in body text if possible, but exact matching might be tricky
                if text and text != slide_data["title"]:
                    slide_data["body_text"].append(text)
        
        # Extract Notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            slide_data["notes"] = slide.notes_slide.notes_text_frame.text.strip()

        content["slides"].append(slide_data)

    with open(output_json_path, "w", encoding="utf-8") as f:
        json.dump(content, f, indent=2, ensure_ascii=False)
    
    print(f"Extraction complete. Saved to {output_json_path}")

if __name__ == "__main__":
    source_ppt = r"C:\Users\Administrator\Downloads\Speeches-and-Statements-of-Quaid-e-Azam-1940-1947.pptx"
    output_json = r"C:\Users\Administrator\.gemini\antigravity\scratch\quaid_content.json"
    extract_content(source_ppt, output_json)
