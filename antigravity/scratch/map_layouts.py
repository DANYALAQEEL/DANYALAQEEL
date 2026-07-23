from pptx import Presentation

path_target = r"c:\Users\Administrator\Downloads\Lecture_Grp_D.pptx"
path_source = r"c:\Users\Administrator\Downloads\Operating_System_Design_Merged.pptx"

print("--- TARGET (Theme) Layouts ---")
prs_target = Presentation(path_target)
for i, layout in enumerate(prs_target.slide_layouts):
    print(f"Index {i}: {layout.name}")

print("\n--- SOURCE (Content) Layouts ---")
prs_source = Presentation(path_source)
for i, slide in enumerate(prs_source.slides):
    layout = slide.slide_layout
    print(f"Slide {i+1}: Uses layout '{layout.name}'")
