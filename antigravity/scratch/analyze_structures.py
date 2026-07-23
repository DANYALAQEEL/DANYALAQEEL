from pptx import Presentation
import os

path_counter = r"c:\Users\Administrator\Downloads\GroupD-SectionA-Topic-2.7.1-2.7.3 (1).pptx"
path_user = r"c:\Users\Administrator\Downloads\Operating_System_Design_Merged.pptx"

def analyze_pptx(path, label):
    if not os.path.exists(path):
        print(f"File not found: {path}")
        return []
        
    prs = Presentation(path)
    print(f"\n--- {label} Structure ({len(prs.slides)} slides) ---")
    slides_data = []
    for i, slide in enumerate(prs.slides):
        title = slide.shapes.title.text if slide.shapes.title else "(No Title)"
        print(f"{i+1}. {title}")
        
        # Capture content summary for later matching
        content = ""
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                content += shape.text_frame.text[:50] + "..."
        slides_data.append({"index": i, "title": title, "content": content})
    return slides_data

print("Analyzing presentations...")
user_slides = analyze_pptx(path_user, "USER")
counter_slides = analyze_pptx(path_counter, "COUNTER GROUP")
