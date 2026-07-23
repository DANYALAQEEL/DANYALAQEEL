from pptx import Presentation
import os
import copy

path_target = r"c:\Users\Administrator\Downloads\Lecture_Grp_D.pptx"
path_target_working = r"c:\Users\Administrator\Downloads\Lecture_Grp_D_Working.pptx"
path_source = r"c:\Users\Administrator\Downloads\Operating_System_Design_Merged.pptx"
output_path = r"c:\Users\Administrator\Downloads\Operating_System_Design_Engineering_Final.pptx"

# Create a working copy of the template
if os.path.exists(path_target_working):
    os.remove(path_target_working)
    
# We will use a different approach:
# 1. Load Target.
# 2. Iterate source slides.
# 3. Append new slides to Target with Target's layout.
# 4. Copy content from Source to these new slides.
# 5. Delete the OLD slides from Target (at the beginning).

print("Loading presentations...")
prs_target = Presentation(path_target)
prs_source = Presentation(path_source)

# Store original slide count to delete later
original_slide_count = len(prs_target.slides)

# Map Source Layout Names to Target Layout Indices
# Based on visual inspection of Lecture_Grp_D
LAYOUT_MAP = {
    "Title Slide": 0,       # TITLE
    "Title and Content": 2, # TITLE_AND_BODY
    "Section Header": 1,    # SECTION_HEADER
    "Two Content": 3,       # TITLE_AND_TWO_COLUMNS (Approximate)
    "Comparison": 3,        # TITLE_AND_TWO_COLUMNS (Approximate)
    "Title Only": 4,        # TITLE_ONLY
    "Blank": 10,            # BLANK
}

def get_target_layout(source_slide):
    name = source_slide.slide_layout.name
    idx = LAYOUT_MAP.get(name, 2) # Default to Title and Content
    # Boundary check
    if idx >= len(prs_target.slide_layouts):
        idx = 2
    return prs_target.slide_layouts[idx]

print("Migrating content...")

for i, source_slide in enumerate(prs_source.slides):
    # 1. Create new slide in target
    layout = get_target_layout(source_slide)
    target_slide = prs_target.slides.add_slide(layout)
    
    # 2. Copy Title
    if source_slide.shapes.title and target_slide.shapes.title:
        target_slide.shapes.title.text = source_slide.shapes.title.text
        
    # 3. Copy Content (Text & Images)
    # Strategy: 
    # - If text placeholder exists in source, move text to target placeholder.
    # - If image exists in source, copy it to target.
    
    # Text Migration
    # Find main body placeholder in source (usually idx 1)
    source_body = None
    for shape in source_slide.placeholders:
        if shape.placeholder_format.idx == 1:
            source_body = shape
            break
            
    if source_body and source_body.has_text_frame:
        # Find target body
        target_body = None
        for shape in target_slide.placeholders:
            if shape.placeholder_format.idx == 1:
                target_body = shape
                break
        
        if target_body:
            # Clear existing text in target placeholder (if any default text)
            target_body.text_frame.clear() 
            
            # Copy paragraphs
            for p in source_body.text_frame.paragraphs:
                new_p = target_body.text_frame.add_paragraph()
                new_p.text = p.text
                new_p.level = p.level
                # Setup font to match theme (don't override unless necessary)
                # But sometimes we might want to keep bold/italic
                if p.font.bold: new_p.font.bold = True
                if p.font.italic: new_p.font.italic = True
                
    # Image Migration (limited support in python-pptx without extraction)
    # We can copy pictures if we can access the blob.
    for shape in source_slide.shapes:
        if shape.shape_type == 13: # PICTURE
            # Add picture to target slide
            # We need to save it to a temp file first
            try:
                image = shape.image
                image_bytes = image.blob
                filename = f"temp_img_{i}_{shape.shape_id}.{image.ext}"
                with open(filename, 'wb') as f:
                    f.write(image_bytes)
                
                # Add to target slide at same position
                target_slide.shapes.add_picture(filename, shape.left, shape.top, shape.width, shape.height)
                os.remove(filename) # Cleanup
            except Exception as e:
                print(f"Failed to copy image on slide {i+1}: {e}")

# Delete original slides from template
# Iterate backwards to avoid index shifting issues
print(f"Deleting {original_slide_count} original template slides...")
# python-pptx doesn't support 'del prs.slides[i]' directly locally?
# Actually it works in newer versions or via xml.
# Standard way involves rIds.
# Let's try the XML method which is robust.
xml_slides = prs_target.slides._sldIdLst
for i in range(original_slide_count):
    # Always delete the FIRST element 'original_slide_count' times? 
    # No, we appended new slides at the END.
    # So the first 'original_slide_count' slides are the old ones.
    # We create a list of the first N elements to remove.
    # Actually, removing the first element N times is safer.
    if len(xml_slides) > 0:
        xml_slides.remove(xml_slides[0])

print(f"Saving to {output_path}...")
prs_target.save(output_path)
print("Done.")
