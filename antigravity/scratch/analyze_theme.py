from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

file_path = r"c:\Users\Administrator\Downloads\Lecture_Grp_D.pptx"
prs = Presentation(file_path)

print(f"Analyzing Theme from: {file_path}")

def get_rgb(color):
    if hasattr(color, 'rgb'):
        return str(color.rgb)
    elif hasattr(color, 'theme_color'):
        return f"Theme Color {color.theme_color}"
    return "None"

# Analyze Slide Master / Layouts
print("\n--- Slide Layouts ---")
for layout in prs.slide_layouts:
    print(f"Layout: {layout.name}")
    bg = layout.background
    if bg:
        fill = bg.fill
        print(f"  Background Fill Type: {fill.type}")
        try:
            print(f"  Background Color: {get_rgb(fill.fore_color)}")
        except:
            pass

# Analyze first few slides for fonts/colors
print("\n--- First 3 Slides Content Analysis ---")
for i, slide in enumerate(prs.slides[:3]):
    print(f"Slide {i+1}")
    
    # Background check on specific slide
    bg = slide.background
    if bg and bg.fill.type:
         print(f"  Custom Background: {bg.fill.type}")
         try:
            print(f"  Custom BG Color: {get_rgb(bg.fill.fore_color)}")
         except: pass

    if slide.shapes.title:
        title = slide.shapes.title
        if title.text_frame.paragraphs:
            p = title.text_frame.paragraphs[0]
            font = p.font
            print(f"  Title Font: {font.name}, Size: {font.size}, Color: {get_rgb(font.color)}")

    # Check other text
    for shape in slide.shapes:
        if shape.has_text_frame and shape != slide.shapes.title:
            for p in shape.text_frame.paragraphs:
                if p.text.strip():
                    font = p.font
                    print(f"  Body Font: {font.name}, Size: {font.size}, Color: {get_rgb(font.color)}")
                    break # Just check the first valid paragraph
