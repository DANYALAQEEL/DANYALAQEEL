from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ────────────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x0D, 0x1B, 0x2A)
GOLD      = RGBColor(0xC8, 0xA8, 0x4B)
TEAL      = RGBColor(0x2A, 0x7F, 0x7F)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE  = RGBColor(0xE8, 0xE0, 0xD0)
DARK_GOLD = RGBColor(0x8B, 0x6D, 0x1F)
LIGHT_NAV = RGBColor(0x16, 0x2A, 0x41)
CREAM     = RGBColor(0xF5, 0xF0, 0xE8)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

# ── Helper functions ────────────────────────────────────────────────────────────
def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)

def bg(slide, color=NAVY):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_w=Pt(1)):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             font_name="Garamond", font_size=Pt(14), bold=False, italic=False,
             color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_para(tf, text, font_name="Garamond", font_size=Pt(13),
             bold=False, italic=False, color=OFFWHITE, align=PP_ALIGN.LEFT, space_before=Pt(4)):
    from pptx.util import Pt as PT
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p

def gold_rule(slide, y, x_start=Inches(0.5), x_end=Inches(12.83)):
    """Thin horizontal gold rule"""
    line = slide.shapes.add_shape(1, x_start, y, x_end - x_start, Pt(1.5))
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()
    return line

def header_bar(slide, title, subtitle=None):
    """Top decorative header bar"""
    # Dark bar
    add_rect(slide, 0, 0, W, Inches(1.3), fill_color=LIGHT_NAV)
    gold_rule(slide, Inches(1.28))
    # Title text
    add_text(slide, title,
             Inches(0.55), Inches(0.1), Inches(11), Inches(0.75),
             font_name="Palatino Linotype", font_size=Pt(30), bold=True,
             color=GOLD, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.55), Inches(0.82), Inches(11), Inches(0.38),
                 font_name="Garamond", font_size=Pt(14), italic=True,
                 color=OFFWHITE, align=PP_ALIGN.LEFT)
    # Corner ornament — small gold square top-right
    add_rect(slide, Inches(12.9), 0, Inches(0.43), Inches(1.3), fill_color=GOLD)

def footer_bar(slide, text="Early Political Leadership of Pakistan  ·  1947–1958"):
    add_rect(slide, 0, Inches(7.13), W, Inches(0.37), fill_color=LIGHT_NAV)
    gold_rule(slide, Inches(7.13))
    add_text(slide, text,
             Inches(0.55), Inches(7.15), Inches(12), Inches(0.3),
             font_name="Garamond", font_size=Pt(9), italic=True,
             color=RGBColor(0xAA, 0xA0, 0x80), align=PP_ALIGN.LEFT)

def side_accent(slide, x=Inches(0.25), color=GOLD, h_frac=0.6):
    """Vertical thin gold bar on left side"""
    bar_h = H * h_frac
    bar_y = (H - bar_h) / 2
    add_rect(slide, x, bar_y, Pt(3), bar_h, fill_color=color)

def tenure_badge(slide, text, x, y):
    """Styled tenure badge"""
    add_rect(slide, x, y, Inches(2.4), Inches(0.38), fill_color=TEAL)
    add_text(slide, text, x + Inches(0.1), y + Pt(2), Inches(2.2), Inches(0.34),
             font_name="Garamond", font_size=Pt(13), bold=True,
             color=WHITE, align=PP_ALIGN.LEFT)

def section_card(slide, title, body, x, y, w=Inches(3.8), h=Inches(1.7),
                 title_color=GOLD, border_color=GOLD):
    """A bordered card with a title and body text"""
    add_rect(slide, x, y, w, h, fill_color=LIGHT_NAV, line_color=border_color, line_w=Pt(1.5))
    # Title strip
    add_rect(slide, x, y, w, Inches(0.38), fill_color=border_color)
    add_text(slide, title, x + Inches(0.1), y + Pt(3), w - Inches(0.15), Inches(0.35),
             font_name="Garamond", font_size=Pt(12), bold=True,
             color=NAVY, align=PP_ALIGN.LEFT)
    txBox = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.42), w - Inches(0.25), h - Inches(0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = body
    run.font.name = "Garamond"
    run.font.size = Pt(11.5)
    run.font.color.rgb = OFFWHITE
    return txBox

def bullet_box(slide, items, x, y, w, h, bullet="▸"):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.name = "Garamond"
        run.font.size = Pt(13)
        run.font.color.rgb = OFFWHITE


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
bg(sl)

# Top thick gold band
add_rect(sl, 0, 0, W, Inches(0.18), fill_color=GOLD)
# Bottom thick gold band
add_rect(sl, 0, Inches(7.32), W, Inches(0.18), fill_color=GOLD)

# Left decorative column
add_rect(sl, 0, 0, Inches(0.55), H, fill_color=LIGHT_NAV)
add_rect(sl, Inches(0.55), 0, Pt(3), H, fill_color=GOLD)

# Right decorative column
add_rect(sl, Inches(12.78), 0, Pt(3), H, fill_color=GOLD)
add_rect(sl, Inches(12.83), 0, Inches(0.5), H, fill_color=LIGHT_NAV)

# Center divider rule
gold_rule(sl, Inches(4.2), x_start=Inches(1.0), x_end=Inches(12.33))

# Main title
add_text(sl, "Early Political Leadership",
         Inches(1.1), Inches(1.2), Inches(11.2), Inches(1.2),
         font_name="Palatino Linotype", font_size=Pt(52), bold=True,
         color=GOLD, align=PP_ALIGN.CENTER)
add_text(sl, "of Pakistan",
         Inches(1.1), Inches(2.25), Inches(11.2), Inches(1.0),
         font_name="Palatino Linotype", font_size=Pt(44), bold=False,
         color=GOLD, align=PP_ALIGN.CENTER)

# Year badge
add_rect(sl, Inches(5.5), Inches(3.3), Inches(2.33), Inches(0.55), fill_color=TEAL)
add_text(sl, "1947 – 1958",
         Inches(5.5), Inches(3.33), Inches(2.33), Inches(0.5),
         font_name="Palatino Linotype", font_size=Pt(18), bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)

# Subtitle
add_text(sl, "Prime Ministers · Governors-General · Constitutional Struggles",
         Inches(1.1), Inches(4.1), Inches(11.2), Inches(0.55),
         font_name="Garamond", font_size=Pt(17), italic=True,
         color=OFFWHITE, align=PP_ALIGN.CENTER)

# Opening quote
add_text(sl,
         '"A nation is born in the hearts of the people long before it takes shape in the laws they make."',
         Inches(2.0), Inches(5.0), Inches(9.33), Inches(0.9),
         font_name="Garamond", font_size=Pt(13), italic=True,
         color=RGBColor(0xAA, 0xA0, 0x80), align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Liaquat Ali Khan
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
bg(sl)
header_bar(sl, "Liaquat Ali Khan", "First Prime Minister of Pakistan")
footer_bar(sl)
side_accent(sl)
tenure_badge(sl, "Tenure: August 1947 – October 1951", Inches(9.3), Inches(1.38))

# Bio paragraph
add_text(sl,
         "Liaquat Ali Khan steered the world's newest nation through its most turbulent infancy. "
         "As Jinnah's right hand during the independence movement, he inherited an almost impossible task: "
         "building a functioning state from scratch while absorbing millions of refugees, managing a partition-era "
         "financial dispute with India over assets, and confronting the Kashmir crisis within weeks of independence.",
         Inches(0.55), Inches(1.42), Inches(7.9), Inches(1.5),
         font_name="Garamond", font_size=Pt(12.5), color=OFFWHITE)

# Key achievements cards
section_card(sl, "Objectives Resolution (1949)",
             "Laid the constitutional foundation of Pakistan as an Islamic state with guaranteed fundamental rights, "
             "free judiciary, and provincial autonomy. Still a preamble to every constitution since.",
             Inches(0.55), Inches(3.05), w=Inches(3.9), h=Inches(1.75))

section_card(sl, "Liaquat–Nehru Pact (1950)",
             "Negotiated a bilateral agreement with India guaranteeing minority rights on both sides of the border, "
             "temporarily calming post-partition communal violence.",
             Inches(4.65), Inches(3.05), w=Inches(3.9), h=Inches(1.75))

section_card(sl, "Assassination — Rawalpindi (1951)",
             "Shot dead at Company Bagh on 16 October 1951. The assassin was immediately killed by police, "
             "and the conspiracy was never fully investigated—leaving a vacuum that destabilised Pakistani democracy permanently.",
             Inches(8.75), Inches(3.05), w=Inches(4.03), h=Inches(1.75), border_color=RGBColor(0xAA, 0x33, 0x33))

# Side quote
add_text(sl, '"The Quaid-e-Millat"',
         Inches(0.55), Inches(2.85), Inches(7.9), Inches(0.35),
         font_name="Garamond", font_size=Pt(12), italic=True,
         color=TEAL, align=PP_ALIGN.LEFT)

# Legacy bar
add_rect(sl, Inches(0.55), Inches(5.05), Inches(12.28), Inches(0.5), fill_color=LIGHT_NAV,
         line_color=GOLD, line_w=Pt(1))
add_text(sl, "Legacy:  Established governance norms, the Objectives Resolution, and Pakistan's first foreign alliances — "
             "his sudden death left democracy without its strongest democratic anchor.",
         Inches(0.7), Inches(5.07), Inches(12.0), Inches(0.44),
         font_name="Garamond", font_size=Pt(11), italic=True, color=OFFWHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Khawaja Nazimuddin
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
bg(sl)
header_bar(sl, "Khawaja Nazimuddin", "2nd Prime Minister · From Prestige to Dismissal")
footer_bar(sl)
side_accent(sl)
tenure_badge(sl, "PM: 1951–1953  |  GG: 1947–1951", Inches(9.0), Inches(1.38))

add_text(sl,
         "Having served as Governor-General after Jinnah's death, Nazimuddin stepped into the premiership after Liaquat's assassination. "
         "A Bengali politician by background, he faced crises on multiple fronts simultaneously — language, religion, and an overreaching Governor-General.",
         Inches(0.55), Inches(1.42), Inches(12.28), Inches(1.1),
         font_name="Garamond", font_size=Pt(12.5), color=OFFWHITE)

# Three-step timeline
steps = [
    ("1951", "Dual Transition",
     "Stepped down as GG to become PM. Held both roles briefly — a sign of already-weakening institutional separation."),
    ("1952", "Language Crisis",
     "His insistence on Urdu as the sole national language triggered the East Pakistan Language Movement. Dhaka students were killed by police on 21 February 1952 — later commemorated as International Mother Language Day."),
    ("1953", "Unconstitutional Dismissal",
     "Governor-General Ghulam Muhammad dismissed him despite a clear parliamentary majority — an act later deemed unconstitutional. It set a precedent that unelected officials could override elected governments."),
]

colors = [TEAL, GOLD, RGBColor(0xAA, 0x33, 0x33)]
for i, (year, title, body) in enumerate(steps):
    x = Inches(0.55) + i * Inches(4.2)
    y = Inches(2.7)
    # Number circle (rectangle approximation)
    add_rect(sl, x, y, Inches(0.55), Inches(0.55), fill_color=colors[i])
    add_text(sl, str(i + 1), x, y, Inches(0.55), Inches(0.55),
             font_name="Palatino Linotype", font_size=Pt(20), bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, f"{year}  ·  {title}",
             x + Inches(0.65), y + Pt(4), Inches(3.35), Inches(0.5),
             font_name="Garamond", font_size=Pt(13.5), bold=True, color=GOLD)
    add_text(sl, body,
             x, y + Inches(0.65), Inches(3.95), Inches(2.4),
             font_name="Garamond", font_size=Pt(11.5), color=OFFWHITE)

# Callout box
add_rect(sl, Inches(0.55), Inches(5.8), Inches(12.28), Inches(0.95),
         fill_color=LIGHT_NAV, line_color=GOLD, line_w=Pt(1))
add_text(sl, "⚖  Democratic Erosion:  Nazimuddin's dismissal proved that in early Pakistan, "
             "the Governor-General's office had become more powerful than Parliament itself — "
             "a structural flaw that would haunt every subsequent government.",
         Inches(0.7), Inches(5.85), Inches(12.0), Inches(0.85),
         font_name="Garamond", font_size=Pt(11.5), italic=True, color=OFFWHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Muhammad Ali Bogra
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
bg(sl)
header_bar(sl, "Muhammad Ali Bogra", "3rd Prime Minister · The Ambassador Who Became PM")
footer_bar(sl)
side_accent(sl)
tenure_badge(sl, "Tenure: April 1953 – August 1955", Inches(9.0), Inches(1.38))

add_text(sl,
         "A surprise appointment lifted directly from his post as Pakistan's Ambassador to the United States, "
         "Bogra had little independent political base — which made him easy for Governor-General Ghulam Muhammad to control. "
         "Despite this, his tenure produced meaningful constitutional and foreign policy milestones.",
         Inches(0.55), Inches(1.42), Inches(12.28), Inches(1.1),
         font_name="Garamond", font_size=Pt(12.5), color=OFFWHITE)

section_card(sl, "Western Alignment",
             "Joined SEATO (1954) and the Baghdad Pact / CENTO (1955). Secured substantial US military and economic aid "
             "during the Cold War, anchoring Pakistan firmly in the Western bloc — a decision with long-lasting consequences.",
             Inches(0.55), Inches(2.75), w=Inches(4.0), h=Inches(2.0))

section_card(sl, "The Bogra Formula (1953)",
             "Proposed a bicameral legislature with equal provincial representation in the upper house and population-based "
             "seats in the lower house — an elegant compromise between East and West Pakistan. The Constituent Assembly "
             "accepted it, but it was never implemented.",
             Inches(4.75), Inches(2.75), w=Inches(4.0), h=Inches(2.0), border_color=TEAL)

section_card(sl, "Puppet Government",
             "When Bogra tried to assert independence, Ghulam Muhammad dissolved the Constituent Assembly in October 1954 "
             "— declared illegal by the Sindh Chief Court, but enforced anyway by sheer executive force. "
             "Bogra was forced to stay on with a new 'cabinet of talents.'",
             Inches(8.95), Inches(2.75), w=Inches(3.93), h=Inches(2.0), border_color=RGBColor(0xAA, 0x33, 0x33))

add_rect(sl, Inches(0.55), Inches(5.0), Inches(12.28), Inches(0.38), fill_color=TEAL)
add_text(sl, "One Unit Scheme (1955):  Bogra's era saw the controversial merging of all West Pakistan provinces into "
             "a single administrative unit — ostensibly to balance East Pakistan's population majority, "
             "but widely resented by Baloch, Pashtun, and Sindhi communities as cultural erasure.",
         Inches(0.7), Inches(5.03), Inches(12.0), Inches(0.32),
         font_name="Garamond", font_size=Pt(11), color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — The Bogra Formula
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
bg(sl)
header_bar(sl, "The Bogra Formula", "A Constitutional Compromise — Built to be Buried")
footer_bar(sl)
side_accent(sl)

add_text(sl,
         "Proposed in October 1953, the Bogra Formula was the most promising attempt to resolve the deadlock "
         "between East Pakistan (majority population) and West Pakistan (political-military dominance). "
         "It proposed a bicameral parliament with a delicate balancing mechanism.",
         Inches(0.55), Inches(1.42), Inches(12.28), Inches(0.95),
         font_name="Garamond", font_size=Pt(12.5), color=OFFWHITE)

# Three pillars diagram
pillars = [
    ("UPPER HOUSE", "80 seats", "Equal representation — 10 seats per unit regardless of population.\nProtected smaller provinces from being overwhelmed.", GOLD),
    ("LOWER HOUSE", "310 seats", "Population-based — 165 for East Pakistan, 145 for West Pakistan.\nGave East Pakistan fair democratic weight for first time.", TEAL),
    ("JOINT SESSION", "Deciding Vote", "For deadlocked bills, both houses sat together. "
     "East Pakistan's majority would prevail — a true concession by West Pakistan.", RGBColor(0x5A, 0x8A, 0x5A)),
]

for i, (title, sub, body, col) in enumerate(pillars):
    x = Inches(0.6) + i * Inches(4.2)
    y = Inches(2.55)
    add_rect(sl, x, y, Inches(3.9), Inches(0.45), fill_color=col)
    add_text(sl, title, x, y + Pt(2), Inches(3.9), Inches(0.42),
             font_name="Palatino Linotype", font_size=Pt(14), bold=True,
             color=NAVY, align=PP_ALIGN.CENTER)
    add_rect(sl, x, y + Inches(0.45), Inches(3.9), Inches(2.55),
             fill_color=LIGHT_NAV, line_color=col, line_w=Pt(1.5))
    add_text(sl, sub, x, y + Inches(0.5), Inches(3.9), Inches(0.45),
             font_name="Garamond", font_size=Pt(15), bold=True,
             color=col, align=PP_ALIGN.CENTER)
    add_text(sl, body, x + Inches(0.15), y + Inches(0.98), Inches(3.6), Inches(1.9),
             font_name="Garamond", font_size=Pt(11.5), color=OFFWHITE)

# Why it failed
add_rect(sl, Inches(0.55), Inches(5.65), Inches(12.28), Inches(1.0),
         fill_color=RGBColor(0x3A, 0x10, 0x10), line_color=RGBColor(0xAA, 0x33, 0x33), line_w=Pt(1.5))
add_text(sl, "✕  Why It Failed:",
         Inches(0.7), Inches(5.68), Inches(3.0), Inches(0.38),
         font_name="Garamond", font_size=Pt(13), bold=True, color=RGBColor(0xFF, 0x66, 0x66))
add_text(sl,
         "Governor-General Ghulam Muhammad dissolved the Constituent Assembly in October 1954 before the formula could be "
         "enacted into law. The courts declared the dissolution illegal — but it stood. Pakistan lost its best chance "
         "at a workable constitutional structure for years.",
         Inches(3.5), Inches(5.68), Inches(9.0), Inches(0.9),
         font_name="Garamond", font_size=Pt(11.5), color=OFFWHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Chaudhry Muhammad Ali
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
bg(sl)
header_bar(sl, "Chaudhry Muhammad Ali", "4th Prime Minister · Father of the 1956 Constitution")
footer_bar(sl)
side_accent(sl)
tenure_badge(sl, "Tenure: August 1955 – September 1956", Inches(8.8), Inches(1.38))

add_text(sl,
         "A career civil servant who had served as Finance Minister and Secretary-General, Chaudhry Muhammad Ali "
         "brought bureaucratic discipline to a chaotic political scene. His singular achievement — passing the 1956 Constitution "
         "after nine years of failed attempts — was also short-lived, as political intrigue forced his resignation within a year.",
         Inches(0.55), Inches(1.42), Inches(12.28), Inches(1.2),
         font_name="Garamond", font_size=Pt(12.5), color=OFFWHITE)

# Constitution details
add_rect(sl, Inches(0.55), Inches(2.75), Inches(7.6), Inches(3.35),
         fill_color=LIGHT_NAV, line_color=GOLD, line_w=Pt(1.5))
add_rect(sl, Inches(0.55), Inches(2.75), Inches(7.6), Inches(0.42), fill_color=GOLD)
add_text(sl, "The 1956 Constitution — Key Features",
         Inches(0.65), Inches(2.77), Inches(7.4), Inches(0.38),
         font_name="Garamond", font_size=Pt(13), bold=True, color=NAVY)
features = [
    "Pakistan formally declared an Islamic Republic for the first time",
    "Parliamentary democracy with Westminster-style cabinet government",
    "Unicameral National Assembly of 156 seats (equal East/West split)",
    "Fundamental rights guaranteed: equality, free speech, religion",
    "Independent judiciary with Federal Court as apex court",
    "Urdu and Bengali recognised as joint official languages",
]
bullet_box(sl, features, Inches(0.7), Inches(3.25), Inches(7.2), Inches(2.7), bullet="◆")

# Right-side cards
section_card(sl, "One Unit Controversy",
             "The merger of West Pakistan into a single province was politically expedient but culturally destructive. "
             "Baloch, Pashtun and Sindhi identities were subsumed — sowing seeds of ethnic resentment that erupted in later decades.",
             Inches(8.35), Inches(2.75), w=Inches(4.48), h=Inches(1.6))

section_card(sl, "Why He Resigned",
             "Crippled by coalition defections orchestrated by President Iskander Mirza and the Republican Party. "
             "Resigned in September 1956 — just months after the Constitution he authored came into force.",
             Inches(8.35), Inches(4.55), w=Inches(4.48), h=Inches(1.55), border_color=RGBColor(0xAA, 0x33, 0x33))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Iskander Mirza
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
bg(sl)
header_bar(sl, "Iskander Mirza", "1st President · The Architect of Pakistan's First Martial Law")
footer_bar(sl)
side_accent(sl, color=RGBColor(0x8B, 0x22, 0x22))

add_text(sl,
         "A former Indian Political Service officer with little faith in parliamentary democracy, Mirza rose through the "
         "ranks of Pakistan's bureaucracy to become the last Governor-General and then the first President. "
         "His political career was defined by intrigue, manipulation, and ultimately the destruction of the system he was meant to protect.",
         Inches(0.55), Inches(1.42), Inches(12.28), Inches(1.1),
         font_name="Garamond", font_size=Pt(12.5), color=OFFWHITE)

# Descent timeline
events = [
    ("1955", "Last Governor-General",
     "Appointed to the office vacated by Ghulam Muhammad. Immediately began using executive power to shape political outcomes.", GOLD),
    ("1956", "First President",
     "Under the new constitution he became a ceremonial president — but refused to act as one. Retained real power through patronage and intimidation.", TEAL),
    ("1956–58", "Political Chess",
     "Cycled through four prime ministers in two years: Chaudhry Muhammad Ali, Suhrawardy, Noon, and Firoz Khan Noon. Created the Republican Party to fragment Muslim League dominance.", RGBColor(0xCC, 0x88, 0x00)),
    ("Oct 1958", "Martial Law",
     "Abrogated the 1956 Constitution, dismissed Noon's government, suspended all political activity, and named General Ayub Khan as Chief Martial Law Administrator.", RGBColor(0xAA, 0x22, 0x22)),
]

for i, (year, title, body, col) in enumerate(events):
    x = Inches(0.55) + i * Inches(3.18)
    y = Inches(2.65)
    add_rect(sl, x, y, Inches(2.9), Inches(0.38), fill_color=col)
    add_text(sl, f"{year}  ·  {title}", x + Pt(5), y + Pt(2), Inches(2.8), Inches(0.35),
             font_name="Garamond", font_size=Pt(11), bold=True, color=NAVY)
    add_rect(sl, x, y + Inches(0.38), Inches(2.9), Inches(2.1),
             fill_color=LIGHT_NAV, line_color=col, line_w=Pt(1))
    add_text(sl, body, x + Pt(8), y + Inches(0.48), Inches(2.75), Inches(1.95),
             font_name="Garamond", font_size=Pt(11), color=OFFWHITE)

# Irony callout
add_rect(sl, Inches(0.55), Inches(5.6), Inches(12.28), Inches(1.1),
         fill_color=LIGHT_NAV, line_color=RGBColor(0xAA, 0x22, 0x22), line_w=Pt(1.5))
add_text(sl, "The Final Irony:",
         Inches(0.7), Inches(5.63), Inches(3), Inches(0.38),
         font_name="Garamond", font_size=Pt(13), bold=True, color=RGBColor(0xFF, 0x66, 0x66))
add_text(sl,
         "Mirza believed he could control Ayub Khan as he had controlled his prime ministers. He was wrong. "
         "Just 20 days after imposing martial law, Ayub Khan exiled Mirza to London — where he died in obscure poverty in 1969. "
         "The man who ended Pakistan's democracy was himself consumed by the forces he had unleashed.",
         Inches(3.6), Inches(5.63), Inches(9.0), Inches(1.0),
         font_name="Garamond", font_size=Pt(11.5), color=OFFWHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Suhrawardy
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
bg(sl)
header_bar(sl, "H.S. Suhrawardy", "5th Prime Minister · The Democrat Who Was Outmanoeuvred")
footer_bar(sl)
side_accent(sl, color=TEAL)
tenure_badge(sl, "Tenure: September 1956 – October 1957", Inches(9.0), Inches(1.38))

add_text(sl,
         "Huseyn Shaheed Suhrawardy was arguably the most capable parliamentarian of the era — a barrister, "
         "Awami League founder, and former Chief Minister of Bengal. He believed sincerely in parliamentary democracy "
         "and fought to make it work against enormous institutional resistance.",
         Inches(0.55), Inches(1.42), Inches(12.28), Inches(1.1),
         font_name="Garamond", font_size=Pt(12.5), color=OFFWHITE)

# Left column — positions
add_rect(sl, Inches(0.55), Inches(2.68), Inches(5.8), Inches(3.3),
         fill_color=LIGHT_NAV, line_color=TEAL, line_w=Pt(1.5))
add_rect(sl, Inches(0.55), Inches(2.68), Inches(5.8), Inches(0.42), fill_color=TEAL)
add_text(sl, "Democratic Positions",
         Inches(0.65), Inches(2.7), Inches(5.6), Inches(0.38),
         font_name="Garamond", font_size=Pt(13), bold=True, color=NAVY)
positions = [
    "Championed genuine parliamentary democracy over presidential excess",
    "Advocated for joint electorates — equal voting across all communities",
    "Pushed for provincial autonomy within a federal structure",
    "Formed the Republican-Awami League coalition — a fragile national unity experiment",
    "Attempted to rebalance the military-civilian power equation",
]
bullet_box(sl, positions, Inches(0.7), Inches(3.18), Inches(5.5), Inches(2.65), bullet="◆")

# Right column — foreign policy + fall
section_card(sl, "Suez Crisis (1956)",
             "In a notable break with Muslim world sentiment, Suhrawardy sided with the Western position during the Suez Canal crisis — "
             "reflecting Pakistan's Cold War alliances but alienating public opinion at home and abroad.",
             Inches(6.6), Inches(2.68), w=Inches(6.23), h=Inches(1.55))

section_card(sl, "The Fall",
             "President Mirza and the Republican Party progressively withdrew support. "
             "Facing an engineered no-confidence threat, Suhrawardy resigned in October 1957 "
             "after just 13 months — denied the stability to execute any long-term democratic programme.",
             Inches(6.6), Inches(4.4), w=Inches(6.23), h=Inches(1.65), border_color=RGBColor(0xAA, 0x33, 0x33))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Malik Firoz Khan Noon
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
bg(sl)
header_bar(sl, "Malik Firoz Khan Noon", "6th & Last PM · Democracy's Final Stand")
footer_bar(sl)
side_accent(sl, color=RGBColor(0xAA, 0x33, 0x33))
tenure_badge(sl, "Tenure: December 1957 – October 1958", Inches(9.0), Inches(1.38))

add_text(sl,
         "A veteran Punjabi politician and former Chief Minister of Punjab, Firoz Khan Noon represented the last genuine "
         "attempt to hold parliamentary governance together. His tenure was short but not without achievement — "
         "and it was ended not by electoral defeat, but by a coup that prevented scheduled elections from ever taking place.",
         Inches(0.55), Inches(1.42), Inches(12.28), Inches(1.15),
         font_name="Garamond", font_size=Pt(12.5), color=OFFWHITE)

# Two-column layout
# Achievements
add_rect(sl, Inches(0.55), Inches(2.72), Inches(5.9), Inches(3.1),
         fill_color=LIGHT_NAV, line_color=GOLD, line_w=Pt(1.5))
add_rect(sl, Inches(0.55), Inches(2.72), Inches(5.9), Inches(0.42), fill_color=GOLD)
add_text(sl, "Achievements",
         Inches(0.65), Inches(2.74), Inches(5.7), Inches(0.38),
         font_name="Garamond", font_size=Pt(13), bold=True, color=NAVY)
achievements = [
    "Negotiated acquisition of Gwadar from the Sultanate of Oman for £3 million (1958) — arguably the most strategically significant act of his tenure",
    "Attempted to stabilise the fractious coalition and proceed toward elections",
    "Maintained cordial foreign relations with both blocs",
]
bullet_box(sl, achievements, Inches(0.7), Inches(3.22), Inches(5.6), Inches(2.45), bullet="◆")

# Crises
add_rect(sl, Inches(6.65), Inches(2.72), Inches(6.08), Inches(3.1),
         fill_color=LIGHT_NAV, line_color=RGBColor(0xAA, 0x33, 0x33), line_w=Pt(1.5))
add_rect(sl, Inches(6.65), Inches(2.72), Inches(6.08), Inches(0.42), fill_color=RGBColor(0xAA, 0x33, 0x33))
add_text(sl, "Crises & Collapse",
         Inches(6.75), Inches(2.74), Inches(5.88), Inches(0.38),
         font_name="Garamond", font_size=Pt(13), bold=True, color=WHITE)
crises = [
    "Governed without a settled parliamentary majority — entirely dependent on shifting coalition allegiances",
    "President Mirza continued to interfere in day-to-day governance",
    "Elections were scheduled for early 1959 — Mirza dismissed his government days before preparations began",
    "Martial Law of October 7, 1958 brought his tenure — and parliamentary democracy — to an abrupt end",
]
bullet_box(sl, crises, Inches(6.8), Inches(3.22), Inches(5.8), Inches(2.45), bullet="▸")

# Footnote
add_text(sl, "★  Gwadar, acquired under Noon, is today the cornerstone of CPEC — perhaps the most consequential legacy of Pakistan's entire first decade of rule.",
         Inches(0.55), Inches(5.98), Inches(12.28), Inches(0.38),
         font_name="Garamond", font_size=Pt(11), italic=True, color=TEAL)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Conclusion
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
bg(sl)
header_bar(sl, "Conclusion", "The Foundation of Political Instability — 11 Years, 6 Prime Ministers")
footer_bar(sl)

# 3 pillars
pillars = [
    ("Executive Dominance",
     "Governors-General and Presidents consistently overrode elected governments. "
     "Between 1947 and 1958, no Prime Minister completed a full term. "
     "The unelected executive held de facto veto power over Parliament.",
     GOLD),
    ("Constitutional Delays",
     "Pakistan took nine years to produce a constitution — longer than any major post-colonial state. "
     "That constitution was abrogated in October 1958, just two years after adoption, "
     "demonstrating how fragile institutional frameworks were.",
     TEAL),
    ("Leadership Instability",
     "Six prime ministers in eleven years meant no policy could take root, no institution could consolidate, "
     "and no democratic norms could form. Each change was driven by presidential manipulation, not democratic mandate.",
     RGBColor(0x5A, 0x8A, 0x5A)),
]

for i, (title, body, col) in enumerate(pillars):
    x = Inches(0.45) + i * Inches(4.25)
    y = Inches(1.5)
    add_rect(sl, x, y, Inches(3.95), Inches(3.4),
             fill_color=LIGHT_NAV, line_color=col, line_w=Pt(1.5))
    add_rect(sl, x, y, Inches(3.95), Inches(0.42), fill_color=col)
    add_text(sl, title, x + Pt(7), y + Pt(3), Inches(3.8), Inches(0.38),
             font_name="Garamond", font_size=Pt(13), bold=True, color=NAVY)
    add_text(sl, body, x + Pt(10), y + Inches(0.5), Inches(3.75), Inches(2.8),
             font_name="Garamond", font_size=Pt(11.5), color=OFFWHITE)

# PM timeline bar
add_rect(sl, Inches(0.45), Inches(5.1), Inches(12.43), Inches(0.38), fill_color=LIGHT_NAV)
gold_rule(sl, Inches(5.1), x_start=Inches(0.45), x_end=Inches(12.88))
gold_rule(sl, Inches(5.46), x_start=Inches(0.45), x_end=Inches(12.88))
pms = ["Liaquat\n1947–51", "Nazimuddin\n1951–53", "Bogra\n1953–55",
       "Ch. M. Ali\n1955–56", "Suhrawardy\n1956–57", "Noon\n1957–58"]
for i, pm in enumerate(pms):
    x = Inches(0.55) + i * Inches(2.05)
    add_text(sl, pm, x, Inches(5.12), Inches(1.95), Inches(0.36),
             font_name="Garamond", font_size=Pt(9), color=OFFWHITE, align=PP_ALIGN.CENTER)

# Closing quote
add_rect(sl, Inches(0.7), Inches(5.7), Inches(11.93), Inches(1.4),
         fill_color=LIGHT_NAV, line_color=GOLD, line_w=Pt(1))
add_text(sl,
         '"The early years of Pakistan\'s independence were marked by a fundamental tension between democratic aspirations '
         'and authoritarian realities — a tension that defined the nation\'s political trajectory and whose echoes '
         'have never truly faded."\n\n— The Pattern Established, 1947–1958',
         Inches(0.85), Inches(5.75), Inches(11.63), Inches(1.3),
         font_name="Garamond", font_size=Pt(12.5), italic=True, color=OFFWHITE, align=PP_ALIGN.CENTER)


# ── Save ───────────────────────────────────────────────────────────────────────
out_path = r"C:\Users\Administrator\Downloads\Early-Pakistan-Leadership-REDESIGNED.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Slides: {len(prs.slides)}")
