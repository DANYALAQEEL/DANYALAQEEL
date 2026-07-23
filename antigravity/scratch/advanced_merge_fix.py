from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
import copy
import os
import time

# Paths
path_counter = r"c:\Users\Administrator\Downloads\GroupD-SectionA-Topic-2.7.1-2.7.3 (1).pptx"
path_user = r"c:\Users\Administrator\Downloads\Operating_System_Design_Merged.pptx"
path_temp = r"c:\Users\Administrator\Downloads\temp_clean_template.pptx"
output_path = r"c:\Users\Administrator\Downloads\Operating_System_Design_Final_Dark_v2.pptx"

# Colors
BG_COLOR = RGBColor(10, 10, 10)
TEXT_TITLE = RGBColor(0, 255, 65)     
TEXT_BODY = RGBColor(220, 220, 220)

print("Phase 1: safely creating clean template...")
# Load User PPT -> Delete All -> Save -> Reload
# This clears internal indices.
try:
    if os.path.exists(path_temp):
        os.remove(path_temp)
        
    prs_clean = Presentation(path_user)
    xml_slides = prs_clean.slides._sldIdLst
    slides_to_remove = list(xml_slides)
    for s in slides_to_remove:
        xml_slides.remove(s)
    prs_clean.save(path_temp)
    print(f"Clean template saved to {path_temp}")
except Exception as e:
    print(f"Error creating template: {e}")
    exit()

print("Phase 2: analyzing and mapping content...")
prs_counter = Presentation(path_counter)
prs_user = Presentation(path_user)

# Map User Slides
def normalize(text):
    return text.lower().strip().replace(":", "").replace("-", " ")

user_map = {}
for i, slide in enumerate(prs_user.slides):
    title = slide.shapes.title.text if slide.shapes.title else ""
    norm_title = normalize(title)
    if norm_title:
        if norm_title not in user_map: user_map[norm_title] = []
        user_map[norm_title].append(slide)

print("Phase 3: constructing final presentation...")
prs_out = Presentation(path_temp) # Load clean template

def apply_theme(slide):
    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR
    
    # Retro Element (Green Line at top)
    try:
        line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(10), Inches(0.1))
        line.fill.solid()
        line.fill.fore_color.rgb = TEXT_TITLE
        line.line.fill.background()
    except Exception as e:
        print(f"Theme Error: {e}")
    
    # Text
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                p.font.name = "Consolas"
                if shape == slide.shapes.title:
                    p.font.color.rgb = TEXT_TITLE
                    p.font.bold = True
                else:
                    p.font.color.rgb = TEXT_BODY

def copy_slide_content(source_slide, target_prs):
    # Create slide using a generic layout
    layout = target_prs.slide_layouts[1] # Title and Content (usually safe)
    
    new_slide = target_prs.slides.add_slide(layout)
    
    # Copy Title
    if source_slide.shapes.title and new_slide.shapes.title:
        new_slide.shapes.title.text = source_slide.shapes.title.text
        
    # Copy Body
    source_body = None
    for shape in source_slide.placeholders:
        if shape.placeholder_format.idx == 1:
            source_body = shape
            break
            
    if source_body and source_body.has_text_frame:
        target_body = None
        for shape in new_slide.placeholders:
            if shape.placeholder_format.idx == 1:
                target_body = shape
                break
        if target_body:
            target_body.text_frame.clear()
            for p in source_body.text_frame.paragraphs:
                new_p = target_body.text_frame.add_paragraph()
                new_p.text = p.text
                new_p.level = p.level
    
    # Copy Images
    for shape in source_slide.shapes:
        if shape.shape_type == 13: # PICTURE
             try:
                image = shape.image
                filename = f"temp_img_{shape.shape_id}.{image.ext}"
                with open(filename, 'wb') as f:
                    f.write(image.blob)
                new_slide.shapes.add_picture(filename, shape.left, shape.top, shape.width, shape.height)
                os.remove(filename)
             except: pass
             
    apply_theme(new_slide)

# 1. Counter Group Structure
used_user_titles = set()
for c_slide in prs_counter.slides:
    c_title = c_slide.shapes.title.text if c_slide.shapes.title else ""
    norm = normalize(c_title)
    
    # Check match
    found_key = None
    for u_key in user_map:
        if u_key in norm or norm in u_key:
            found_key = u_key
            break
            
    if found_key:
        print(f"Merging Topic (User): {found_key}")
        for u_slide in user_map[found_key]:
            copy_slide_content(u_slide, prs_out)
        used_user_titles.add(found_key)
    else:
        print(f"Merging Topic (Counter): {c_title}")
        copy_slide_content(c_slide, prs_out)

# 2. Append Orphans
for u_key, u_slides in user_map.items():
    if u_key not in used_user_titles:
        print(f"Appending Extra: {u_key}")
        for u_slide in u_slides:
            copy_slide_content(u_slide, prs_out)

print(f"Saving final to {output_path}...")
prs_out.save(output_path)
print("Success.")
