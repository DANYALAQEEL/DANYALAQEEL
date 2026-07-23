from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os
import copy

path_counter = r"c:\Users\Administrator\Downloads\GroupD-SectionA-Topic-2.7.1-2.7.3 (1).pptx"
path_user = r"c:\Users\Administrator\Downloads\Operating_System_Design_Merged.pptx"
output_path = r"c:\Users\Administrator\Downloads\Operating_System_Design_Final_GrpD_Theme.pptx"

print("Loading presentations...")
prs_out = Presentation(path_counter) # Base is Counter Group (Template)
prs_user = Presentation(path_user)

# Map User Slides
def normalize(text):
    return text.lower().strip().replace(":", "").replace("-", " ")

user_map = {}
for i, slide in enumerate(prs_user.slides):
    title = slide.shapes.title.text if slide.shapes.title else ""
    norm_title = normalize(title)
    if norm_title:
        if norm_title not in user_map: user_map[norm_title] = []
        user_map[norm_title].append(slide)

def clear_slide_content(slide):
    # Remove all shapes safely
    list_shapes = list(slide.shapes)
    for shape in list_shapes:
        sp = shape.element
        sp.getparent().remove(sp)

def copy_content_to_slide(source_slide, target_slide):
    # 1. Title
    if source_slide.shapes.title:
        # Check if target has a title placeholder (it should if layout has it)
        if target_slide.shapes.title: 
            target_slide.shapes.title.text = source_slide.shapes.title.text
        else:
            # Create a title box if missing (rare)
            title_shape = target_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
            title_shape.text_frame.text = source_slide.shapes.title.text
        
    # 2. Body
    source_body = None
    for shape in source_slide.placeholders:
        if shape.placeholder_format.idx == 1:
            source_body = shape
            break
    
    # Find Target Body
    target_body = None
    # If layout was preserved, placeholder should exist
    for shape in target_slide.placeholders:
        if shape.placeholder_format.idx == 1:
            target_body = shape
            break
            
    if source_body and source_body.has_text_frame:
        if target_body:
            # Use existing placeholder to keep THEME formatting
            target_body.text_frame.clear()
            tf = target_body.text_frame
        else:
            # Create textbox if no placeholder
            body_shape = target_slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
            tf = body_shape.text_frame

        for p in source_body.text_frame.paragraphs:
            new_p = tf.add_paragraph()
            new_p.text = p.text
            new_p.level = p.level
            # Do NOT copy font/color to preserve Template Theme
            # Only copy bold/italic if essential
            if p.font.bold: new_p.font.bold = True
            if p.font.italic: new_p.font.italic = True
            
    # 3. Images
    for shape in source_slide.shapes:
        if shape.shape_type == 13:
             try:
                image = shape.image
                filename = f"temp_theme_{shape.shape_id}.{image.ext}"
                with open(filename, 'wb') as f:
                    f.write(image.blob)
                target_slide.shapes.add_picture(filename, shape.left, shape.top, shape.width, shape.height)
                os.remove(filename)
             except: pass

def add_new_slide_copy(source_slide, prs):
    # Determine layout
    # Try to map source layout to target layout
    # Simple mapping:
    # Title Slide -> 0
    # Title and Content -> 2
    # Section Header -> 1
    # Title Only -> 4 (usually)
    
    src_name = source_slide.slide_layout.name
    layout_idx = 2 # Default Body
    
    if "Title Slide" in src_name: layout_idx = 0
    elif "Section Header" in src_name: layout_idx = 1
    elif "Title Only" in src_name: layout_idx = 4
    
    if layout_idx >= len(prs.slide_layouts): layout_idx = 2
    
    layout = prs.slide_layouts[layout_idx] 
    new_slide = prs.slides.add_slide(layout)
    copy_content_to_slide(source_slide, new_slide)

print("Merging with Group D Theme...")
used_user_titles = set()

# Iterate through Counter Group Slides (Base)
for i, c_slide in enumerate(prs_out.slides):
    c_title = c_slide.shapes.title.text if c_slide.shapes.title else ""
    norm = normalize(c_title)
    
    found_key = None
    for u_key in user_map:
        if u_key in norm or norm in u_key:
            found_key = u_key
            break
            
    if found_key:
        print(f"Replacing Topic (User): {found_key}")
        # Clear content but KEEP LAYOUT and BACKGROUND of the template
        # We need to be careful not to delete the placeholders if we want to use them!
        # clear_slide_content DELETES shapes. This removes placeholders.
        # Instead, we should CLEAR PLACEHOLDERS.
        
        # Safe Clear Strategy:
        # 1. Clear Title Text
        if c_slide.shapes.title: c_slide.shapes.title.text_frame.clear()
        
        # 2. Clear Body Text
        for shape in c_slide.placeholders:
            if shape.placeholder_format.idx == 1 and shape.has_text_frame:
                shape.text_frame.clear()
        
        # 3. Remove other shapes (images/textboxes that are not placeholders)
        # This is harder to distinguish safely.
        # For now, let's just Append into the placeholders.
        
        u_slides = user_map[found_key]
        copy_content_to_slide(u_slides[0], c_slide)
        used_user_titles.add(found_key)
    else:
        print(f"Keeping Topic (Counter): {c_title}")
        # Automatically keeps theme since it IS the file.

# Append Orphans
for u_key, u_slides in user_map.items():
    if u_key not in used_user_titles:
        print(f"Appending Extra: {u_key}")
        for u_slide in u_slides:
            add_new_slide_copy(u_slide, prs_out)

print(f"Saving to {output_path}")
prs_out.save(output_path)
print("Done.")
