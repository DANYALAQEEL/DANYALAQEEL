import pptx
import os

def analyze_presentation(pptx_path):
    if not os.path.exists(pptx_path):
        print(f"Error: File not found at {pptx_path}")
        return

    prs = pptx.Presentation(pptx_path)
    
    print(f"Presentation: {os.path.basename(pptx_path)}")
    print(f"Total Slides: {len(prs.slides)}")
    print("-" * 30)

    slide_data = []

    for i, slide in enumerate(prs.slides):
        slide_content = {
            "slide_number": i + 1,
            "title": "",
            "text": []
        }
        
        if slide.shapes.title:
            slide_content["title"] = slide.shapes.title.text.strip()
            
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text = shape.text.strip()
                if text and text != slide_content["title"]:
                     slide_content["text"].append(text)
        
        slide_data.append(slide_content)
        
        print(f"Slide {i+1}: {slide_content['title']}")
        if slide_content['text']:
            for text_item in slide_content['text'][:3]: # Print first 3 text items to avoid clutter
                print(f"  - {text_item[:100]}...")
            if len(slide_content['text']) > 3:
                print(f"  ... and {len(slide_content['text']) - 3} more items")
        print("-" * 30)

if __name__ == "__main__":
    ppt_path = "Operating_System_Design_Engineering_Final_Organized.pptx"
    analyze_presentation(ppt_path)
