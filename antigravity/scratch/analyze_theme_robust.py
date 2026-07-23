from pptx import Presentation

file_path = r"c:\Users\Administrator\Downloads\Lecture_Grp_D.pptx"
try:
    prs = Presentation(file_path)
except Exception as e:
    print(f"Error loading presentation: {e}")
    exit()

print(f"Successfully loaded: {file_path}")

# Check Slide Master
try:
    master = prs.slide_masters[0]
    print("\n--- Slide Master ---")
    
    # Background
    bg = master.background
    if bg:
        print(f"Master Background Element: {bg.element.xml[:100]}...") # Print a bit of XML to see what we have
    
    # Theme Fonts
    tm = master.text_styles
    if tm:
        print("Text Styles found in master.")
        
except Exception as e:
    print(f"Error analyzing master: {e}")

# Check Layouts (safely)
print("\n--- Layouts ---")
for i, layout in enumerate(prs.slide_layouts):
    print(f"Layout {i}: {layout.name}")
    try:
        if layout.slide_master:
            print(f"  Linked to Master: {layout.slide_master.name}")
    except:
        print("  Could not access master for this layout.")

# Try to get *one* slide safely
try:
    if len(prs.slides) > 0:
        slide = prs.slides[0]
        print("\n--- First Slide Analysis ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                print(f"  Shape text: {shape.text[:50]}")
                # Try to get font info
                try:
                    p = shape.text_frame.paragraphs[0]
                    # Direct check on font
                    print(f"    Font: {p.font.name}, Size: {p.font.size}")
                    # Color check
                    if p.font.color:
                        if hasattr(p.font.color, 'rgb'):
                            print(f"    Color RGB: {p.font.color.rgb}")
                        if hasattr(p.font.color, 'theme_color'):
                            print(f"    Color Theme: {p.font.color.theme_color}")
                except:
                    pass
except Exception as e:
    print(f"Error accessing first slide: {e}")
