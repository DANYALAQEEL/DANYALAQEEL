from pptx import Presentation
import os

path_target = r"c:\Users\Administrator\Downloads\Lecture_Grp_D.pptx"
path_source = r"c:\Users\Administrator\Downloads\Operating_System_Design_Merged.pptx"
output_path = r"c:\Users\Administrator\Downloads\Operating_System_Design_Engineering_Final.pptx"

print("Loading presentations...")
prs_target = Presentation(path_target) # Template
prs_source = Presentation(path_source) # Content

# Function to get layout from target based on source layout
def get_target_layout(source_layout_name, target_prs):
    # Mapping logic based on previous analysis
    name_map = {
        "Title Slide": 0,       # TITLE
        "Title and Content": 2, # TITLE_AND_BODY
        "Section Header": 1,    # SECTION_HEADER
        "Two Content": 3,       # TITLE_AND_TWO_COLUMNS
        "Comparison": 3,        # TITLE_AND_TWO_COLUMNS
        "Title Only": 4,        # TITLE_ONLY
        "Blank": 6,             # BLANK (actually index 10 in target, but let's check)
    }
    
    # Fallback to index 2 (Title and Content) if unknown
    idx = name_map.get(source_layout_name, 2)
    
    # Specific fix for blank if needed, previous analysis showed BLANK at index 10
    if source_layout_name == "Blank": idx = 10
    
    if idx < len(target_prs.slide_layouts):
        return target_prs.slide_layouts[idx]
    else:
        return target_prs.slide_layouts[2]

# Remove existing slides from target to make it a clean template
# Note: python-pptx doesn't have a simple "delete slide"
# We have to use xml manipulation or just create a new one based on master?
# Actually, easier way: iterating backwards and removing from xml
def delete_all_slides(prs):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    for s in slides:
        xml_slides.remove(s)

print("Clearing template slides...")
delete_all_slides(prs_target)

print("Migrating content...")
for i, source_slide in enumerate(prs_source.slides):
    # 1. Determine Layout
    layout = get_target_layout(source_slide.slide_layout.name, prs_target)
    
    # 2. Create New Slide
    new_slide = prs_target.slides.add_slide(layout)
    
    # 3. Copy Title
    if source_slide.shapes.title and new_slide.shapes.title:
        new_slide.shapes.title.text = source_slide.shapes.title.text
        
    # 4. Copy Body Text
    # We look for the main placeholder in source and target
    # Usually index 1 for body
    source_body = None
    for shape in source_slide.placeholders:
        if shape.placeholder_format.idx == 1:
            source_body = shape
            break
            
    if source_body and source_body.has_text_frame:
        # Find target body placeholder
        target_body = None
        for shape in new_slide.placeholders:
            if shape.placeholder_format.idx == 1:
                target_body = shape
                break
        
        if target_body:
            target_body.text = source_body.text_frame.text
            # Basic formatting transfer could happen here if needed, 
            # but we want to KEEP target formatting mostly.
            
    # 5. Copy Images (basic ext step)
    # This is complex in pptx without extracting files. 
    # For now, we will notify if we skipped images or implement a rigorous extraction if requested.
    # Given the user wants "engineering satisfying theme", text content is priority.
    
print(f"Saving to {output_path}...")
prs_target.save(output_path)
print("Done.")
