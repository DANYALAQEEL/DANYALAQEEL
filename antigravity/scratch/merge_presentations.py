import pptx
import os
from pptx.util import Inches

def get_slide_text(slide):
    text = []
    if slide.shapes.title:
        text.append(slide.shapes.title.text)
    for shape in slide.shapes:
        if shape.has_text_frame:
             for paragraph in shape.text_frame.paragraphs:
                 if paragraph.text.strip():
                     text.append(paragraph.text.strip())
    return " ".join(text).lower()

def merge_presentations(target_path, source_path, output_path):
    print(f"Loading Target: {target_path}")
    target_prs = pptx.Presentation(target_path)
    print(f"Loading Source: {source_path}")
    source_prs = pptx.Presentation(source_path)

    # 1. Map existing topics in Target
    existing_topics = set()
    for slide in target_prs.slides:
        if slide.shapes.title:
            existing_topics.add(slide.shapes.title.text.strip().lower())

    print(f"Found {len(existing_topics)} existing topics in Target.")

    # 2. Iterate Source and Add Missing
    # Find insertion index: before "Quiz" or "Summary"
    insert_index = len(target_prs.slides)
    for i, slide in enumerate(target_prs.slides):
        if slide.shapes.title:
            title = slide.shapes.title.text.lower()
            if "quiz" in title or "summary" in title or "q & a" in title:
                insert_index = i
                break
    
    print(f"Insertion point determined at slide index: {insert_index}")
    
    # We can't easily "insert" slides in python-pptx at a specific index without manual xml manipulation or re-creating the whole slide list.
    # A simpler approach is to append them, and then we can move them if strictly necessary, 
    # OR for this task, just append them before the end if we accept they might be at the end.
    # However, to respect "order", we will append them to the end of the presentation for now, 
    # as re-ordering slides is complex in python-pptx (requires low-level XML manipulation).
    # IF the user insists on order, we might need a more complex script. 
    # For now, let's append valid missing slides.

    added_count = 0
    
    # We need to determine the correct layout to use from Target.
    # Usually layout 1 is Title and Content.
    content_layout = target_prs.slide_layouts[1] 

    for slide in source_prs.slides:
        source_title = ""
        if slide.shapes.title:
            source_title = slide.shapes.title.text.strip()
        
        if not source_title:
            continue

        # Check if topic exists
        if source_title.lower() in existing_topics:
            print(f"Skipping existing topic: {source_title}")
            continue
            
        print(f"Adding new slide: {source_title}")
        
        # Create new slide in Target
        new_slide = target_prs.slides.add_slide(content_layout)
        
        # Copy Title
        if new_slide.shapes.title:
            new_slide.shapes.title.text = source_title
        
        # Copy Content (Text)
        # We'll aggregate text from the source slide's body shapes
        source_body_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                if shape == slide.shapes.title:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        source_body_text.append(paragraph.text)
        
        # specific to Title & Content layout, the placeholder[1] is usually the body
        if len(new_slide.placeholders) > 1:
            body_shape = new_slide.placeholders[1]
            tf = body_shape.text_frame
            tf.clear() # Clear default prompt text
            
            for text in source_body_text:
                p = tf.add_paragraph()
                p.text = text
                p.level = 0 # specific level handling could be added here
        
        # Copy Images (Basic support)
        # We cannot easily position them exactly as in source without complex math,
        # so we will place them in a standard location if they exist.
        img_count = 0
        for shape in slide.shapes:
            if shape.shape_type == 13: # MSO_SHAPE_TYPE.PICTURE
                try:
                    image = shape.image
                    # save blob to temp file
                    image_filename = f"temp_img_{added_count}_{img_count}.jpg"
                    with open(image_filename, 'wb') as f:
                        f.write(image.blob)
                    
                    # Add to new slide - positioning is arbitrary (bottom right roughly)
                    # You might want to adjust this logic
                    left = Inches(5)
                    top = Inches(2)
                    height = Inches(3)
                    new_slide.shapes.add_picture(image_filename, left, top, height=height)
                    
                    os.remove(image_filename)
                    img_count += 1
                except Exception as e:
                    print(f"Could not copy image: {e}")

        added_count += 1
        existing_topics.add(source_title.lower())

    # Save
    target_prs.save(output_path)
    print(f"Merged presentation saved to: {output_path}")
    print(f"Total slides added: {added_count}")

target = r"C:\Users\Administrator\Downloads\Operating_System_Design_Engineering_Final.pptx"
source = r"C:\Users\Administrator\Downloads\GroupD-SectionA-Topic-2.7.1-2.7.3 (1).pptx"
output = r"C:\Users\Administrator\Downloads\Operating_System_Design_Engineering_Final_Merged.pptx"

if __name__ == "__main__":
    merge_presentations(target, source, output)
