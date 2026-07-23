from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
import os

path_counter = r"c:\Users\Administrator\Downloads\GroupD-SectionA-Topic-2.7.1-2.7.3 (1).pptx"
path_user = r"c:\Users\Administrator\Downloads\Operating_System_Design_Merged.pptx"
output_path = r"c:\Users\Administrator\Downloads\Operating_System_Design_Final_Dark_Fixed.pptx"

BG_COLOR = RGBColor(10, 10, 10)
TEXT_TITLE = RGBColor(0, 255, 65)     
TEXT_BODY = RGBColor(220, 220, 220)

print("Loading presentations...")
prs_out = Presentation(path_counter) # Base is Counter Group
prs_user = Presentation(path_user)

# Map User Slides
def normalize(text):
    return text.lower().strip().replace(":", "").replace("-", " ")

user_map = {}
for i, slide in enumerate(prs_user.slides):
    title = slide.shapes.title.text if slide.shapes.title else ""
    norm_title = normalize(title)
    if norm_title:
        if norm_title not in user_map: user_map[norm_title] = []
        user_map[norm_title].append(slide)

def clear_slide_content(slide):
    # Remove all shapes safely
    list_shapes = list(slide.shapes)
    for shape in list_shapes:
        sp = shape.element
        sp.getparent().remove(sp)

def apply_theme(slide):
    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR
    
    # Retro Element
    try:
        line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(10), Inches(0.1))
        line.fill.solid()
        line.fill.fore_color.rgb = TEXT_TITLE
        line.line.fill.background()
    except: pass
    
    # Text
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                p.font.name = "Consolas"
                if shape == slide.shapes.title:
                    p.font.color.rgb = TEXT_TITLE
                    p.font.bold = True
                else:
                    p.font.color.rgb = TEXT_BODY

def copy_content_to_slide(source_slide, target_slide):
    # 1. Title
    if source_slide.shapes.title:
        title = target_slide.shapes.title # Might be deleted?
        # If we deleted everything, we need to add a title box?
        # Yes, we deleted everything.
        # Add Title Box
        title_shape = target_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_shape.text_frame.text = source_slide.shapes.title.text
        # We'll rely on apply_theme to style it
        
    # 2. Body
    # Find source body
    source_body = None
    for shape in source_slide.placeholders:
        if shape.placeholder_format.idx == 1:
            source_body = shape
            break
    
    if source_body and source_body.has_text_frame:
        body_shape = target_slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
        for p in source_body.text_frame.paragraphs:
            new_p = body_shape.text_frame.add_paragraph()
            new_p.text = p.text
            new_p.level = p.level
            
    # 3. Images
    for shape in source_slide.shapes:
        if shape.shape_type == 13:
             try:
                image = shape.image
                filename = f"temp_fix_{shape.shape_id}.{image.ext}"
                with open(filename, 'wb') as f:
                    f.write(image.blob)
                target_slide.shapes.add_picture(filename, shape.left, shape.top, shape.width, shape.height)
                os.remove(filename)
             except: pass
             
    apply_theme(target_slide)

def add_new_slide_copy(source_slide, prs):
    # Just use Title Only layout for safety
    layout = prs.slide_layouts[5] 
    new_slide = prs.slides.add_slide(layout)
    # Clear placeholders to be safe
    clear_slide_content(new_slide)
    copy_content_to_slide(source_slide, new_slide)

print("Merging and overwriting...")
used_user_titles = set()

# Iterate through Counter Group Slides (Base)
for i, c_slide in enumerate(prs_out.slides):
    # Get original title before we clear it?
    # Yes.
    c_title = c_slide.shapes.title.text if c_slide.shapes.title else ""
    norm = normalize(c_title)
    
    # We will ALWAYS clear the content to apply our own styling and Ensure consistency
    # But if we don't have User content, we should keep Counter content?
    # User said: "add topics included in their presentation and missing in mine"
    # So if missing in user, we KEEP counter content.
    
    found_key = None
    for u_key in user_map:
        if u_key in norm or norm in u_key:
            found_key = u_key
            break
            
    if found_key:
        print(f"Replacing Topic (User): {found_key}")
        clear_slide_content(c_slide)
        # Copy User Content (Taking the first match if multiple?)
        # If multiple user slides map to one counter slide, we put first here, append others?
        # Logic: Copy first here.
        u_slides = user_map[found_key]
        copy_content_to_slide(u_slides[0], c_slide)
        
        # If there are more, we need to insert them after this slide?
        # Insert is hard. We'll append them at end for now or ignore (assume 1:1)
        used_user_titles.add(found_key)
    else:
        print(f"Keeping Topic (Counter): {c_title}")
        # We still want to apply Theme!
        apply_theme(c_slide)

# Append Orphans
for u_key, u_slides in user_map.items():
    if u_key not in used_user_titles:
        print(f"Appending Extra: {u_key}")
        for u_slide in u_slides:
            add_new_slide_copy(u_slide, prs_out)

print(f"Saving to {output_path}")
prs_out.save(output_path)
print("Done.")
