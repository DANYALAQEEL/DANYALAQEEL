import fitz
import os
from pptx import Presentation

pdf_path = r"C:\Users\Administrator\Downloads\Text Book - Operating System Concepts (10th Ed) - Gagne Silberschatz and Galvin 2018.pdf"
base_dir = r"C:\Users\Administrator\.gemini\antigravity\scratch"

# Define chapter page ranges (1-indexed based on the toc, so we need to subtract 1 for 0-indexed PyMuPDF)
chapters_pages = {
    1: (30, 83),
    2: (84, 141),
    3: (142, 214),
    4: (215, 264),
    5: (265, 330)
}

doc = fitz.open(pdf_path)

def extract_pdf_text():
    for chapter, (start, end) in chapters_pages.items():
        text = ""
        for page_num in range(start, end + 1):
            page = doc.load_page(page_num)
            text += page.get_text() + "\n"
        
        out_path = os.path.join(base_dir, f"chapter_{chapter}_pdf.txt")
        with open(out_path, "w", encoding="utf-8") as f:
            f.write(text)
        print(f"Extracted Chapter {chapter} from PDF to {out_path}")

ppt_paths = {
    1: r"C:\Users\Administrator\Downloads\CS-330-Operating Systems 2K24-BESE-15 Sp26_2026033_0413\19 January - 25 January\Lesson 1- Intro\Lesson 1- Intro.pptx",
    2: r"C:\Users\Administrator\Downloads\CS-330-Operating Systems 2K24-BESE-15 Sp26_2026033_0413\26 January - 1 February\Chapter 2 - Operating-System Services\Chapter 2 - Operating-System Services.pptx",
    3: r"C:\Users\Administrator\Downloads\CS-330-Operating Systems 2K24-BESE-15 Sp26_2026033_0413\9 February - 15 February\Chapter 3 - Processes\Chapter 3 - Processes.pptx",
    4: r"C:\Users\Administrator\Downloads\CS-330-Operating Systems 2K24-BESE-15 Sp26_2026033_0413\16 February - 22 February\Chapter 4 - Threads and Concurrency\Chapter 4 - Threads and Concurrency.pptx",
    5: r"C:\Users\Administrator\Downloads\CS-330-Operating Systems 2K24-BESE-15 Sp26_2026033_0413\23 February - 1 March\Chapter 5 - CPU Scheduling\Chapter 5 - CPU Scheduling.pptx"
}

def extract_ppt_text():
    for chapter, ppt_path in ppt_paths.items():
        if os.path.exists(ppt_path):
            try:
                prs = Presentation(ppt_path)
                text = ""
                for slide in prs.slides:
                    # Extract slide text
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                    # Extract slide notes
                    if slide.has_notes_slide:
                        notes_slide = slide.notes_slide
                        text += "--- NOTES ---\n"
                        text += notes_slide.notes_text_frame.text + "\n"
                    text += "\n==============\n"
                    
                out_path = os.path.join(base_dir, f"chapter_{chapter}_ppt.txt")
                with open(out_path, "w", encoding="utf-8") as f:
                    f.write(text)
                print(f"Extracted Chapter {chapter} from PPT to {out_path}")
            except Exception as e:
                print(f"Failed to process {ppt_path}: {e}")
        else:
            print(f"Could not find PPT for chapter {chapter} at {ppt_path}")

extract_pdf_text()
extract_ppt_text()
