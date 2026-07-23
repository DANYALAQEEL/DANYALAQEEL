from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import math
import matplotlib.pyplot as plt
import numpy as np
import os

def create_presentation():
    prs = Presentation()

    # --- Theme Settings (Slate) ---
    BG_COLOR = RGBColor(40, 44, 52) # Dark Slate
    TEXT_COLOR = RGBColor(220, 223, 228) # Light Grey/White
    ACCENT_COLOR = RGBColor(97, 175, 239) # Light Blue
    
    FORMULAS = {
        "Newton's Forward Difference": "P(x) = y0 + u*dy0 + [u(u-1)/2!]*d2y0 + ...",
        "Newton's Backward Difference": "P(x) = yn + u*dyn + [u(u+1)/2!]*d2yn + ...",
        "Newton's Divided Difference": "P(x) = b0 + b1(x-x0) + b2(x-x0)(x-x1) + ...",
        "Lagrange Interpolation": "P(x) = Sum[yj * Product((x-xi)/(xj-xi))]"
    }

    def apply_slate_theme(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR
        if slide.shapes.title:
            slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR

    def add_slide(layout_index, title_text, content_text):
        slide_layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(slide_layout)
        apply_slate_theme(slide)
        title = slide.shapes.title
        title.text = title_text
        for p in title.text_frame.paragraphs:
            p.font.name = 'Arial'
            p.font.size = Pt(28)
            p.font.color.rgb = ACCENT_COLOR

        if content_text:
            body = slide.shapes.placeholders[1]
            tf = body.text_frame
            tf.text = content_text
            for p in tf.paragraphs:
                p.font.name = 'Arial'
                p.font.size = Pt(16)
                p.font.color.rgb = TEXT_COLOR
        return slide

    def add_table_slide(title_text, data, col_names):
        slide_layout = prs.slide_layouts[5] # Title Only
        slide = prs.slides.add_slide(slide_layout)
        apply_slate_theme(slide)
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
        title.text_frame.paragraphs[0].font.size = Pt(28)

        # Dynamic sizing
        rows = len(data) + 1
        cols = len(col_names)
        
        # Adjust existing table layout logic
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9.0)
        
        # Calculate row height based on count to fit slide
        row_height = 0.5
        if rows > 8: row_height = 0.4
        height = Inches(row_height * rows)

        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Header
        for i, name in enumerate(col_names):
            cell = table.cell(0, i)
            cell.text = str(name)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(60, 64, 72)
            for p in cell.text_frame.paragraphs:
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.font.bold = True
                p.font.size = Pt(12)

        # Data
        for i, row in enumerate(data):
            for j, val in enumerate(row):
                cell = table.cell(i+1, j)
                # Format floats nicely
                if isinstance(val, float):
                    cell.text = f"{val:.4f}"
                else:
                    cell.text = str(val)
                
                cell.fill.solid()
                cell.fill.fore_color.rgb = BG_COLOR
                for p in cell.text_frame.paragraphs:
                    p.font.color.rgb = TEXT_COLOR
                    p.font.size = Pt(12)

    def add_image_slide(title_text, image_path, content_text=""):
        slide_layout = prs.slide_layouts[5] # Title Only
        slide = prs.slides.add_slide(slide_layout)
        apply_slate_theme(slide)
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
        
        top = Inches(1.5)
        left = Inches(0.5)
        height = Inches(5.0)
        slide.shapes.add_picture(image_path, left, top, height=height)
        
        # Add text box for content
        if content_text:
            text_box = slide.shapes.add_textbox(Inches(6.0), top, Inches(3.5), Inches(5.0))
            tf = text_box.text_frame
            tf.text = content_text
            tf.word_wrap = True
            for p in tf.paragraphs:
                p.font.color.rgb = TEXT_COLOR
                p.font.size = Pt(14)

    # --- Math & Graphing ---
    def generate_plot(x, y, target_x, result_y, method_name, scenario_id):
        plt.figure(figsize=(6, 4), facecolor='#282c34')
        ax = plt.gca()
        ax.set_facecolor('#282c34')
        
        # Scatter original points
        plt.scatter(x, y, color='#61afef', label='Data Points', s=100, zorder=3)
        
        # Plot Interpolated Point
        plt.scatter([target_x], [result_y], color='#e06c75', label='Interpolated', s=150, zorder=4, marker='*')

        # Smooth Curve
        try:
            # We'll use a high-degree polynomial fit (numpy) as a visual approximation 
            # for the specific method curve, or Lagrange for generic visualization
            x_smooth = np.linspace(min(x) - 1, max(x) + 1, 200)
            
            # Use Lagrange implementation for the curve to be accurate to the method
            # (Newton and Lagrange generate the same polynomial)
            y_smooth = []
            for val in x_smooth:
                # Reuse lagrange logic for curve generation
                # lagrange returns: result, steps, display_table, cols
                res_val = lagrange(x, y, val)[0]
                y_smooth.append(res_val)
            
            plt.plot(x_smooth, y_smooth, color='#98c379', linestyle='--', label='Curve', alpha=0.7)
        except Exception as e:
            print(f"Plot curve error: {e}")

        plt.title(f"{method_name} (S{scenario_id})", color='white')
        plt.tick_params(colors='white')
        plt.grid(True, color='#3e4451', linestyle=':', alpha=0.6)
        plt.legend()
        
        filename = f"plot_s{scenario_id}_{method_name.replace(' ', '_')}.png"
        plt.savefig(filename, dpi=100, bbox_inches='tight')
        plt.close()
        return filename

    def check_intervals(x):
        if len(x) < 2: return True
        h = round(x[1] - x[0], 5)
        for i in range(1, len(x) - 1):
            if round(x[i+1] - x[i], 5) != h: return False
        return True

    def newton_forward(x, y, target):
        n = len(x)
        h = x[1] - x[0]
        u = (target - x[0]) / h
        diff_table = [y[:]]
        for i in range(1, n):
            prev = diff_table[-1]
            curr = [prev[j+1] - prev[j] for j in range(len(prev)-1)]
            diff_table.append(curr)
        
        # Reformat table for display: Rows = x values, Cols = y, dy, d2y...
        # Standard notation: Row i contains x[i], y[i], dy[i], d2y[i]... 
        # But dy[i] is typically placed between rows. We will just list them straight.
        display_table = []
        cols = ["x", "y"] + [f"D{k}y" for k in range(1, n)]
        
        for i in range(n):
            row_data = [x[i]] # Start with x
            # Append diff terms that exist for this row index
            # diff_table[k] has length n-k.
            # Forward diff: entry [i] in col [k] corresponds to delta^k y_i
            for k in range(n):
                if i < len(diff_table[k]):
                    row_data.append(diff_table[k][i])
                else:
                    row_data.append("")
            display_table.append(row_data)

        result = y[0]
        u_term = 1
        fact = 1
        steps = f"u = {u:.4f}\n"
        for i in range(1, n):
            u_term *= (u - (i-1))
            fact *= i
            term = (u_term * diff_table[i][0]) / fact
            result += term
            steps += f"Term {i}: {term:.4f}\n"
            
        return result, steps, display_table, cols

    def newton_backward(x, y, target):
        n = len(x)
        h = x[1] - x[0]
        u = (target - x[-1]) / h
        diff_table = [y[:]]
        for i in range(1, n):
            prev = diff_table[-1]
            curr = [prev[j+1] - prev[j] for j in range(len(prev)-1)]
            diff_table.append(curr)

        display_table = []
        cols = ["x", "y"] + [f"D{k}y" for k in range(1, n)]
        for i in range(n):
            row_data = [x[i]]
            for k in range(n):
                if i < len(diff_table[k]):
                    row_data.append(diff_table[k][i])
                else:
                    row_data.append("")
            display_table.append(row_data)

        result = y[-1]
        u_term = 1
        fact = 1
        steps = f"u = {u:.4f}\n"
        for i in range(1, n):
            u_term *= (u + (i-1))
            fact *= i
            if not diff_table[i]: break
            val = diff_table[i][-1]
            term = (u_term * val) / fact
            result += term
            steps += f"Term {i}: {term:.4f}\n"

        return result, steps, display_table, cols

    def divided_difference(x, y, target):
        n = len(x)
        table = [[0] * n for _ in range(n)]
        for i in range(n): table[i][0] = y[i]
        for j in range(1, n):
            for i in range(n - j):
                table[i][j] = (table[i+1][j-1] - table[i][j-1]) / (x[i+j] - x[i])
        
        display_table = []
        cols = ["x", "y"] + [f"DD{k}" for k in range(1, n)]
        for i in range(n):
            row_data = [x[i]]
            for j in range(n):
                if j < (n-i):
                   row_data.append(table[i][j])
                else:
                   row_data.append("")
            display_table.append(row_data)

        result = table[0][0]
        prod = 1
        steps = f"b0 = {table[0][0]:.4f}\n"
        for i in range(1, n):
            prod *= (target - x[i-1])
            term = table[0][i] * prod
            result += term
            steps += f"Term {i}: {term:.4f}\n"
        
        return result, steps, display_table, cols

    def lagrange(x, y, target):
        n = len(x)
        result = 0
        steps = ""
        for i in range(n):
            term = y[i]
            for j in range(n):
                if i != j:
                    term *= (target - x[j]) / (x[i] - x[j])
            result += term
            steps += f"L{i}: {term:.4f}\n"
        
        # No diff table for Lagrange, just return None
        return result, steps, None, None

    # --- Data Scenarios ---
    scenarios = [
        {
            "id": 1, "field": "Chem Eng (Thermodynamics)", "problem": "Find Cp at T=305K",
            "x": [300, 320, 340, 360], "y": [1.005, 1.008, 1.013, 1.020], "t": 305,
            "cols": ["Temp (K)", "Cp"]
        },
        {
            "id": 2, "field": "Economics (Population)", "problem": "Find Pop in 2008",
            "x": [1980, 1990, 2000, 2010], "y": [10, 12, 15, 20], "t": 2008,
            "cols": ["Year", "Pop (M)"]
        },
        {
            "id": 3, "field": "Physics (Kinematics)", "problem": "Find v at t=2.5s",
            "x": [0, 1, 3, 4], "y": [0, 10, 22, 28], "t": 2.5,
            "cols": ["Time (s)", "Vel (m/s)"]
        },
        {
            "id": 4, "field": "Robotics (Trajectory)", "problem": "Find y at x=4",
            "x": [1, 2, 5, 7], "y": [3, 5, 12, 8], "t": 4,
            "cols": ["x", "y"]
        }
    ]

    methods = [
        ("Newton's Forward Difference", newton_forward, True),
        ("Newton's Backward Difference", newton_backward, True),
        ("Newton's Divided Difference", divided_difference, False),
        ("Lagrange Interpolation", lagrange, False)
    ]

    # --- Generate Slides ---
    add_slide(0, "Polynomial Interpolation: 16-Case Enhanced", 
              "Analysis with Formulas, Difference Tables, and Graphs.\n\nGenerated by Antigravity")

    for s in scenarios:
        add_slide(1, f"Scenario {s['id']}: {s['field']}",
                  f"Problem: {s['problem']}\n"
                  f"Dataset: x={s['x']}, y={s['y']}\n"
                  f"Target: {s['t']}")
        
        data_rows = [[xv, yv] for xv, yv in zip(s['x'], s['y'])]
        add_table_slide(f"Scenario {s['id']} Data", data_rows, s['cols'])

        is_equal = check_intervals(s['x'])

        for m_name, m_func, req_equal in methods:
            title_pfx = f"S{s['id']} - {m_name}"
            
            if req_equal and not is_equal:
                add_slide(1, title_pfx, 
                          "Method NOT Applicable.\n"
                          "Reason: Unequal intervals.\n"
                          "Formula: " + FORMULAS.get(m_name, ""))
            else:
                try:
                    res, steps, d_table, d_cols = m_func(s['x'], s['y'], s['t'])
                    
                    # 1. Formula & Table Slide
                    if d_table:
                        add_table_slide(f"{title_pfx} (Table)", d_table, d_cols)
                    
                    # 2. Graph & Result Slide
                    plot_file = generate_plot(s['x'], s['y'], s['t'], res, m_name, s['id'])
                    
                    content = (f"Formula: {FORMULAS.get(m_name, '')}\n\n"
                               f"Calculation Steps:\n{steps}\n"
                               f"Final Result: {res:.5f}")
                    
                    add_image_slide(f"{title_pfx} (Result)", plot_file, content)
                    
                    # Cleanup plot file
                    if os.path.exists(plot_file):
                        os.remove(plot_file)
                        
                except Exception as e:
                    add_slide(1, title_pfx, f"Error: {e}")

    prs.save('Interpolation_Examples_Enhanced.pptx')
    print("Saved 'Interpolation_Examples_Enhanced.pptx'")

if __name__ == "__main__":
    create_presentation()
