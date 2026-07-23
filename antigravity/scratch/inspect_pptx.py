from pptx import Presentation
import os

file_path = r"c:\Users\Administrator\Downloads\Operating_System_Design_Merged.pptx"

if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    exit()

prs = Presentation(file_path)

print(f"Presentation: {os.path.basename(file_path)}")
print(f"Total Slides: {len(prs.slides)}")

for i, slide in enumerate(prs.slides):
    print(f"\n--- Slide {i+1} ---")
    if slide.shapes.title:
        print(f"Title: {slide.shapes.title.text}")
    else:
        print("Title: (No Title)")
    
    for shape in slide.shapes:
        print(f"  Shape: {shape.name} (Type: {shape.shape_type})")
        if hasattr(shape, "text") and shape.text:
            print(f"    Text: {shape.text[:50]}...")
