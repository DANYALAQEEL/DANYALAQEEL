import pptx
import os

def analyze_pptx(filepath):
    print(f"--- Analyzing: {os.path.basename(filepath)} ---")
    if not os.path.exists(filepath):
        print(f"Error: File not found at {filepath}")
        return

    try:
        prs = pptx.Presentation(filepath)
        print(f"Total Slides: {len(prs.slides)}")
        for i, slide in enumerate(prs.slides):
            title = ""
            if slide.shapes.title:
                title = slide.shapes.title.text
            # print(f"Slide {i+1}: {title}") # Redundant for just checking

    except Exception as e:
        print(f"Error reading file: {e}")

target_path = r"C:\Users\Administrator\Downloads\Operating_System_Design_Engineering_Final_Organized_Themed_Fixed.pptx"
analyze_pptx(target_path)
