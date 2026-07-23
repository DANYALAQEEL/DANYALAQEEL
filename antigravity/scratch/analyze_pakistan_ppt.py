from pptx import Presentation
from pptx.util import Inches, Pt
import json

pptx_path = r"C:\Users\Administrator\Downloads\Early-Political-Leadership-of-Pakistan-1947-1958.pptx"
prs = Presentation(pptx_path)

slides_content = []

for i, slide in enumerate(prs.slides):
    slide_data = {
        "slide_number": i + 1,
        "layout_name": slide.slide_layout.name if slide.slide_layout else "Unknown",
        "shapes": []
    }
    
    for shape in slide.shapes:
        shape_data = {
            "shape_name": shape.name,
            "shape_type": str(shape.shape_type),
            "has_text": shape.has_text_frame,
        }
        
        if shape.has_text_frame:
            paragraphs = []
            for para in shape.text_frame.paragraphs:
                para_text = para.text.strip()
                if para_text:
                    paragraphs.append(para_text)
            shape_data["text"] = paragraphs
        
        slide_data["shapes"].append(shape_data)
    
    slides_content.append(slide_data)

# Print summary
for slide in slides_content:
    print(f"\n{'='*60}")
    print(f"SLIDE {slide['slide_number']} (Layout: {slide['layout_name']})")
    print('='*60)
    for shape in slide['shapes']:
        if shape.get('text'):
            print(f"  [{shape['shape_name']}]:")
            for line in shape['text']:
                print(f"    - {line}")

# Save to JSON
with open(r"C:\Users\Administrator\.gemini\antigravity\scratch\pakistan_ppt_content.json", "w", encoding='utf-8') as f:
    json.dump(slides_content, f, indent=2, ensure_ascii=False)

print("\n\nTotal slides:", len(slides_content))
print("Content saved to pakistan_ppt_content.json")
