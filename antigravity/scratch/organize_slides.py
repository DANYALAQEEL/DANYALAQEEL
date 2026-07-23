import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

def create_divider_slide(prs, title_text):
    slide_layout = prs.slide_layouts[2] # Section Header
    slide = prs.slides.add_slide(slide_layout)
    if slide.shapes.title:
        slide.shapes.title.text = title_text
    return slide

def organize_presentation(input_path, output_path):
    print(f"Loading: {input_path}")
    prs = pptx.Presentation(input_path)

    # Use slide_id as unique identifier for sets
    slides_by_title = {}
    original_slides = [] # Store tuple (slide_obj, title, slide_id)
    
    for slide in prs.slides:
        title = ""
        if slide.shapes.title:
            title = slide.shapes.title.text.strip()
        original_slides.append((slide, title, slide.slide_id))
    
    # Define Sections
    sections = {
        "INTRODUCTION": [
            "Operating System Design & Implementation",
            "Agenda",
            "Introduction",
            "The Nature of the Problem",
            "OS Design – The Art of the Trade-Off",
            "Operating System Design and Implementation",
        ],
        "DESIGN GOALS & CONFLICTS": [
            "The Two Sets of Goals: The Core Conflict",
            "Resolving the Conflict",
            "User Goals in Mobile Context",
            "System Goals",
        ],
        "CORE CONCEPTS": [
            "The OS as the Ultimate Middleman",
            "Mechanism vs. Policy",
            "Separation in Practice – Domain Examples",
        ],
        "ARCHITECTURES": [
            "Monolithic Structure",
            "Microkernel Structure",
            "Modular and Hybrid Systems",
            "Modular Systems",
            "Hybrid Systems",
        ],
        "IMPLEMENTATION": [
            "Implementation – The Language of the Kernel",
            "Evolution of OS Code",
            "The Modern OS Stack",
            "Why C Dominates",
            "The Performance Myth",
            "Bridge Concept: Emulation",
            "The Modern Shift – The Rust Revolution",
            "System Generation (SYSGEN)",
            "The Boot Process",
        ],
        "ASSESSMENT & CONCLUSION": [
            "Assessment",
            "Quiz: Concepts",
            "Quiz: Architecture",
            "Summary & Takeaways",
            "Key Takeaways",
            "Final Summary",
            "Q & A",
        ]
    }

    used_slide_ids = set()
    xml_slides = prs.slides._sldIdLst
    new_xml_order = []

    # Iterate sections
    for section_name, keywords in sections.items():
        print(f"Processing Section: {section_name}")
        
        # 1. Add Divider Slide - Appends to XML list
        divider = prs.slides.add_slide(prs.slide_layouts[2])
        if divider.shapes.title:
            divider.shapes.title.text = f">>> {section_name} <<<"
        
        # The divider is now the LAST element in xml_slides
        divider_xml = xml_slides[-1]
        new_xml_order.append(divider_xml)

        # 2. Find matching slides
        for kw in keywords:
            for slide, title, s_id in original_slides:
                if s_id in used_slide_ids:
                    continue
                
                if kw.lower() in title.lower():
                    print(f"  Found: {title}")
                    
                    # Find xml element index
                    idx = [s[2] for s in original_slides].index(s_id)
                    slide_xml = xml_slides[idx]
                    
                    if slide_xml not in new_xml_order:
                        new_xml_order.append(slide_xml)
                        used_slide_ids.add(s_id)

    # 3. Add Leftovers
    print("Processing Leftovers...")
    divider = prs.slides.add_slide(prs.slide_layouts[2])
    if divider.shapes.title:
        divider.shapes.title.text = ">>> APPENDIX / OTHER <<<"
    new_xml_order.append(xml_slides[-1])

    for slide, title, s_id in original_slides:
        if s_id not in used_slide_ids:
            print(f"  Leftover: {title}")
            idx = [s[2] for s in original_slides].index(s_id)
            slide_xml = xml_slides[idx]
            if slide_xml not in new_xml_order:
                new_xml_order.append(slide_xml)

    # Apply new order
    print(f"Applying new order. Total slides: {len(new_xml_order)}")
    
    # We must operate on list copy
    original_sldId_list = list(xml_slides)
    
    # Check checks
    if len(new_xml_order) != (len(original_slides) + len(sections) + 1):
        print(f"Warning: Count mismatch! Expected {len(original_slides) + len(sections) + 1} but got {len(new_xml_order)}")

    # Clear current list
    for sldId in original_sldId_list:
        xml_slides.remove(sldId)
    
    # Append new order
    for sldId in new_xml_order:
        xml_slides.append(sldId)
        
    prs.save(output_path)
    print(f"Saved categorized presentation to: {output_path}")

input_path = r"C:\Users\Administrator\Downloads\Operating_System_Design_Engineering_Final_Merged.pptx"
output_path = r"C:\Users\Administrator\Downloads\Operating_System_Design_Engineering_Final_Organized.pptx"

if __name__ == "__main__":
    organize_presentation(input_path, output_path)
